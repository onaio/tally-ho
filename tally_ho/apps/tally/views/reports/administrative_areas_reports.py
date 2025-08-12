import ast
import datetime
import json
import os
from datetime import date
from io import BytesIO

from django.conf import settings
from django.contrib.postgres.aggregates import ArrayAgg
from django.db.models import (
    Case,
    CharField,
    Count,
    ExpressionWrapper,
    F,
    IntegerField,
    OuterRef,
    Q,
    Subquery,
    Sum,
    When,
)
from django.db.models import Value as V
from django.db.models.functions import Coalesce
from django.http import HttpResponse, JsonResponse
from django.shortcuts import redirect
from django.urls import reverse
from django.utils import timezone
from django.utils.translation import gettext_lazy as _
from django.views.generic import TemplateView
from django_datatables_view.base_datatable_view import BaseDatatableView
from guardian.mixins import LoginRequiredMixin
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from tally_ho.apps.tally.models.all_candidates_votes import AllCandidatesVotes
from tally_ho.apps.tally.models.center import Center
from tally_ho.apps.tally.models.constituency import Constituency
from tally_ho.apps.tally.models.electrol_race import ElectrolRace
from tally_ho.apps.tally.models.reconciliation_form import ReconciliationForm
from tally_ho.apps.tally.models.region import Region
from tally_ho.apps.tally.models.result import Result
from tally_ho.apps.tally.models.result_form import ResultForm
from tally_ho.apps.tally.models.station import Station
from tally_ho.apps.tally.models.sub_constituency import SubConstituency
from tally_ho.apps.tally.views.reports.helpers import (
    get_filtered_candidate_votes,
)
from tally_ho.apps.tally.views.super_admin import (
    get_result_form_with_duplicate_results,
)
from tally_ho.libs.models.enums.entry_version import EntryVersion
from tally_ho.libs.models.enums.form_state import FormState
from tally_ho.libs.permissions import groups
from tally_ho.libs.utils.numbers import parse_int
from tally_ho.libs.utils.query_set_helpers import Round
from tally_ho.libs.views.mixins import (
    DataTablesMixin,
    GroupRequiredMixin,
    TallyAccessMixin,
)

report_types = {
    1: "turnout",
    2: "summary",
    3: "stations-and-centers-under-process-audit-list",
    4: "stations-and-centers-under-investigation-list",
    5: "stations-and-centers-excluded-after-investigation-list",
    6: "progressive-report",
}


def build_stations_centers_and_sub_cons_list(tally_id):
    """
    Create a list of stations and centers filtered by tally id.

    :param tally_id: Tally id.

    returns: list of stations, centers and sub cons.
    """
    qs = Station.objects.filter(tally__id=tally_id).distinct(
        "tally__id", "center__code", "station_number"
    )

    stations = list(
        qs.annotate(name=F("station_number"))
        .values("name")
        .annotate(id=F("id"))
    )
    centers = list(
        qs.annotate(name=F("center__code"))
        .values("name")
        .annotate(id=F("center__id"))
        .distinct("center__code")
    )
    sub_cons = list(
        qs.annotate(name=F("center__sub_constituency__name"))
        .values("name")
        .annotate(code=F("center__sub_constituency__code"))
        .distinct("center__sub_constituency__code")
    )

    return stations, centers, sub_cons


def get_stations_and_centers_by_admin_area(
    tally_id,
    report_column_name,
    report_type_name,
    region_id=None,
    constituency_id=None,
):
    """
    Genarate a report of stations and centers under investigation or excluded
    after investigation.

    :param tally_id: The reconciliation forms tally.
    :param report_column_name: The result form report column name.
    :param report_type_name: The report type name to generate.
    :param region_id: The result form report region id used for filtering.
    :param constituency_id: The result form report constituency id
        used for filtering.

    returns: The stations and centers report grouped by the adminstrative
        area name.
    """
    qs = Station.objects.filter(tally__id=tally_id)

    stations_centers_audit_report_type_name = report_types[3]
    stations_centers_under_investigation_report_type_name = report_types[4]
    stations_centers_excluded_after_investigation_report_type_name = (
        report_types[5]
    )
    centers_count_query = None

    if report_type_name == stations_centers_audit_report_type_name:
        qs = qs.filter(active=False)
        centers_count_query = Subquery(
            Center.objects.annotate(
                center_count=Coalesce(
                    Count("id", filter=Q(tally__id=tally_id, active=False)),
                    V(0),
                )
            ).values("center_count")[:1],
            output_field=IntegerField(),
        )

    if (
        report_type_name
        == stations_centers_under_investigation_report_type_name
    ):
        qs = qs.filter(active=False)
        centers_count_query = Subquery(
            Center.objects.annotate(
                center_count=Coalesce(
                    Count("id", filter=Q(tally__id=tally_id, active=False)),
                    V(0),
                )
            ).values("center_count")[:1],
            output_field=IntegerField(),
        )

    if (
        report_type_name
        == stations_centers_excluded_after_investigation_report_type_name
    ):
        qs = qs.filter(
            Q(active=True, center__disable_reason__isnull=False)
            | Q(active=True, disable_reason__isnull=False)
        )
        centers_count_query = Subquery(
            Center.objects.annotate(
                center_count=Coalesce(
                    Count(
                        "id",
                        filter=Q(
                            tally__id=tally_id,
                            active=True,
                            disable_reason__isnull=False,
                        ),
                    ),
                    V(0),
                )
            ).values("center_count")[:1],
            output_field=IntegerField(),
        )

    if region_id:
        qs = qs.filter(center__office__region__id=region_id)
    if constituency_id:
        qs = qs.filter(center__constituency__id=constituency_id)

    qs = (
        qs.annotate(
            admin_area_name=F(report_column_name),
            region_id=F("center__office__region__id"),
        )
        .values(
            "admin_area_name",
            "region_id",
        )
        .annotate(
            number_of_centers=centers_count_query,
            number_of_stations=Count("station_number"),
            total_number_of_centers_and_stations=ExpressionWrapper(
                F("number_of_centers") + F("number_of_stations"),
                output_field=IntegerField(),
            ),
        )
    )

    if region_id:
        qs = qs.annotate(
            constituency_id=F("center__constituency__id"),
            sub_constituency__id=F("sub_constituency__id"),
        )

    return qs


def stations_and_centers_queryset(tally_id, qs, data=None, **kwargs):
    """
    Genarate a report of stations and centers under investigation or excluded
    after investigation.

    :param tally_id: The reconciliation forms tally.
    :param report_column_name: The result form report column name.
    :param report_type_name: The report type name to generate.
    :param region_id: The result form report region id used for filtering.
    :param constituency_id: The result form report constituency id
        used for filtering.

    returns: The stations and centers report grouped by the adminstrative
        area name.
    """
    stations_centers_under_process_audit = report_types[3]
    stations_centers_under_investigation_report_type = report_types[4]
    stations_centers_excluded_after_investigation_report_type = report_types[5]
    column_name = "center__office__region__name"
    column_id = "center__office__region__id"
    admin_area_id = kwargs.get("region_id")
    constituency_id = kwargs.get("constituency_id")
    report_type = kwargs.get("report_type")

    if admin_area_id and constituency_id:
        column_name = "center__sub_constituency__code"
        column_id = "center__sub_constituency__id"
    elif admin_area_id:
        column_id = "center__constituency__id"
        column_name = "center__constituency__name"

    qs = qs.filter(tally__id=tally_id)

    if admin_area_id:
        qs = qs.filter(center__office__region__id=admin_area_id)
    if constituency_id:
        qs = qs.filter(center__constituency__id=constituency_id)

    if data and len(
        [
            data
            for d in data
            if len(d["select_1_ids"]) or len(d["select_2_ids"])
        ]
    ):
        parent_qs = qs
        for item in data:
            item["region_id"]

            center_ids = (
                item["select_2_ids"] if len(item["select_2_ids"]) else [0]
            )
            station_ids = (
                item["select_1_ids"] if len(item["select_1_ids"]) else [0]
            )

            if report_type == stations_centers_under_investigation_report_type:
                # if admin_area_id and constituency_id:
                #     current_qs = parent_qs.filter(
                #         center__sub_constituency__id=region_id
                #     )
                # elif admin_area_id:
                #     current_qs = parent_qs.filter(
                #         center__constituency__id=region_id
                #     )
                # else:
                #     current_qs = parent_qs.filter(
                #         center__office__region__id=region_id
                #     )

                current_qs = (
                    parent_qs.filter(active=False)
                    .annotate(
                        admin_area_name=F(column_name),
                        admin_area_id=F(column_id),
                    )
                    .values(
                        "admin_area_name",
                        "admin_area_id",
                    )
                    .annotate(
                        region_id=F("center__office__region__id"),
                        number_of_centers=Count(
                            "center__id",
                            distinct=True,
                            filter=(
                                ~Q(center__id__in=center_ids)
                                & Q(center__active=False)
                            ),
                        ),
                        number_of_stations=Count(
                            "station_number", filter=(~Q(id__in=station_ids))
                        ),
                        total_number_of_centers_and_stations=ExpressionWrapper(
                            F("number_of_centers") + F("number_of_stations"),
                            output_field=IntegerField(),
                        ),
                        center_ids=ArrayAgg(
                            "center__id",
                            distinct=True,
                            filter=Q(center__active=False),
                        ),
                        station_ids=ArrayAgg("id", distinct=True),
                        constituency_id=F("center__constituency__id"),
                        sub_constituency_id=F("center__sub_constituency__id"),
                    )
                )

                qs = (
                    qs.union(current_qs)
                    if not isinstance(qs[0], Station)
                    else current_qs
                )

            elif (
                report_type
                == stations_centers_excluded_after_investigation_report_type
            ):
                # if admin_area_id and constituency_id:
                #     current_qs = parent_qs.filter(
                #         center__sub_constituency__id=region_id
                #     )
                # elif admin_area_id:
                #     current_qs = parent_qs.filter(
                #         center__constituency__id=region_id
                #     )
                # else:
                #     current_qs = parent_qs.filter(
                #         center__office__region__id=region_id
                #     )

                current_qs = (
                    parent_qs.filter(
                        Q(active=True, center__disable_reason__isnull=False)
                        | Q(active=True, disable_reason__isnull=False)
                    )
                    .annotate(
                        admin_area_name=F(column_name),
                        admin_area_id=F(column_id),
                    )
                    .values(
                        "admin_area_name",
                        "admin_area_id",
                    )
                    .annotate(
                        region_id=F("center__office__region__id"),
                        number_of_centers=Count(
                            "center__id",
                            distinct=True,
                            filter=(
                                ~Q(center__id__in=center_ids)
                                & Q(
                                    center__active=True,
                                    center__disable_reason__isnull=False,
                                )
                            ),
                        ),
                        number_of_stations=Count(
                            "station_number", filter=(~Q(id__in=station_ids))
                        ),
                        total_number_of_centers_and_stations=ExpressionWrapper(
                            F("number_of_centers") + F("number_of_stations"),
                            output_field=IntegerField(),
                        ),
                        center_ids=ArrayAgg(
                            "center__id",
                            distinct=True,
                            filter=(
                                Q(
                                    center__active=True,
                                    center__disable_reason__isnull=False,
                                )
                            ),
                        ),
                        station_ids=ArrayAgg("id", distinct=True),
                        constituency_id=F("center__constituency__id"),
                        sub_constituency_id=F("center__sub_constituency__id"),
                    )
                )

                qs = (
                    qs.union(current_qs)
                    if not isinstance(qs[0], Station)
                    else current_qs
                )

            elif report_type == stations_centers_under_process_audit:
                station_numbers = (
                    list(
                        Station.objects.filter(
                            tally__id=tally_id, id__in=station_ids
                        ).values_list("station_number", flat=True)
                    )
                    if not (len(station_ids) == 1 and not station_ids[0])
                    else [0]
                )

                # if admin_area_id and constituency_id:
                #     current_qs = parent_qs.filter(
                #         center__sub_constituency__id=region_id
                #     )
                # elif admin_area_id:
                #     current_qs = parent_qs.filter(
                #         center__constituency__id=region_id
                #     )
                # else:
                #     current_qs = parent_qs.filter(
                #         center__office__region__id=region_id
                #     )

                current_qs = (
                    parent_qs.annotate(
                        admin_area_name=F(column_name),
                        admin_area_id=F(column_id),
                    )
                    .values(
                        "admin_area_name",
                        "admin_area_id",
                    )
                    .annotate(
                        region_id=F("center__office__region__id"),
                        number_of_centers=Count(
                            "center__id",
                            distinct=True,
                            filter=~Q(center__id__in=center_ids),
                        ),
                        number_of_stations=Count(
                            "station_number",
                            filter=(~Q(station_number__in=station_numbers)),
                        ),
                        total_number_of_centers_and_stations=ExpressionWrapper(
                            F("number_of_centers") + F("number_of_stations"),
                            output_field=IntegerField(),
                        ),
                        center_ids=ArrayAgg("center__id", distinct=True),
                        station_ids=ArrayAgg("station_number", distinct=True),
                    )
                )
                if admin_area_id:
                    current_qs = current_qs.annotate(
                        constituency_id=F("center__constituency__id"),
                        sub_constituency_id=F("center__sub_constituency__id"),
                    )

                qs = (
                    qs.union(current_qs)
                    if not isinstance(qs[0], ResultForm)
                    else current_qs
                )
    else:
        qs = qs.annotate(
            admin_area_name=F(column_name), admin_area_id=F(column_id)
        ).values(
            "admin_area_name",
            "admin_area_id",
        )

        if report_type == stations_centers_under_investigation_report_type:
            qs = qs.filter(active=False).annotate(
                number_of_centers=Count(
                    "center__id", distinct=True, filter=Q(center__active=False)
                ),
                station_ids=ArrayAgg("id", distinct=True),
                center_ids=ArrayAgg(
                    "center__id",
                    distinct=True,
                    filter=(
                        Q(
                            center__active=False,
                            center__disable_reason__isnull=False,
                        )
                    ),
                ),
                constituency_id=F("center__constituency__id"),
                sub_constituency_id=F("center__sub_constituency__id"),
            )
        elif (
            report_type
            == stations_centers_excluded_after_investigation_report_type
        ):
            qs = qs.filter(
                Q(active=True, center__disable_reason__isnull=False)
                | Q(active=True, disable_reason__isnull=False)
            ).annotate(
                number_of_centers=Count(
                    "center__id",
                    distinct=True,
                    filter=Q(
                        center__active=True,
                        center__disable_reason__isnull=False,
                    ),
                ),
                station_ids=ArrayAgg("id", distinct=True),
                center_ids=ArrayAgg(
                    "center__id",
                    distinct=True,
                    filter=(
                        Q(
                            center__active=True,
                            center__disable_reason__isnull=False,
                        )
                    ),
                ),
                constituency_id=F("center__constituency__id"),
                sub_constituency_id=F("center__sub_constituency__id"),
            )
        elif report_type == stations_centers_under_process_audit:
            qs = qs.annotate(
                number_of_centers=Count("center__id", distinct=True),
                station_ids=ArrayAgg("station_number", distinct=True),
                center_ids=ArrayAgg("center__id", distinct=True),
            )

            if admin_area_id:
                qs = qs.annotate(
                    constituency_id=F("center__constituency__id"),
                    sub_constituency_id=F("center__sub_constituency__id"),
                )

        qs = qs.annotate(
            number_of_stations=Count("station_number", distinct=True),
            total_number_of_centers_and_stations=ExpressionWrapper(
                F("number_of_centers") + F("number_of_stations"),
                output_field=IntegerField(),
            ),
            region_id=F("center__office__region__id"),
        )

    return qs


def get_total_valid_votes_per_electrol_race(tally_id, electral_race_id):
    total_valid_votes = (
        Result.objects.filter(
            result_form__tally__id=tally_id,
            active=True,
            result_form__ballot__electrol_race__id=electral_race_id,
            entry_version=EntryVersion.FINAL,
            result_form__form_state=FormState.ARCHIVED,
        )
        .values("result_form__ballot__electrol_race__id")
        .aggregate(
            total_valid_votes=Sum("votes"),
        )
        .get("total_valid_votes")
    )
    return total_valid_votes or 0


def results_queryset(tally_id, qs, data=None):
    """
    Genarate a report of votes per candidate.

    :param tally_id: The tally id.
    :param qs: The result parent queryset.
    :param data: An array of dicts containing centers and stations
        id's to filter out from the queryset.

    returns: The votes per candidate queryset.
    """
    station_id_query = Subquery(
        Station.objects.filter(
            tally__id=tally_id,
            center__code=OuterRef("result_form__center__code"),
            station_number=OuterRef("result_form__station_number"),
        ).values("id")[:1],
        output_field=IntegerField(),
    )

    station_total_result_forms_sub_query = Subquery(
        ResultForm.objects.filter(
            tally__id=tally_id,
            center__code=OuterRef("center__code"),
            station_number=OuterRef("station_number"),
        )
        .values("center__code", "station_number")
        .annotate(
            total_result_forms=Count("barcode", distinct=True),
        )
        .values("total_result_forms")[:1],
        output_field=IntegerField(),
    )

    station_total_result_forms_archived_sub_query = Subquery(
        ResultForm.objects.filter(
            tally__id=tally_id,
            center__code=OuterRef("center__code"),
            station_number=OuterRef("station_number"),
        )
        .values("center__code", "station_number")
        .annotate(
            total_result_forms_archived=Count(
                "barcode",
                distinct=True,
                filter=Q(form_state=FormState.ARCHIVED),
            ),
        )
        .values("total_result_forms_archived")[:1],
        output_field=IntegerField(),
    )

    if data:
        selected_center_ids = (
            data["select_1_ids"] if data.get("select_1_ids") else []
        )
        selected_station_ids = (
            data["select_2_ids"] if data.get("select_2_ids") else []
        )
        election_level_names = (
            data["election_level_names"]
            if data.get("election_level_names")
            else []
        )
        sub_race_type_names = (
            data["sub_race_type_names"]
            if data.get("sub_race_type_names")
            else []
        )
        ballot_status = (
            data["ballot_status"] if data.get("ballot_status") else []
        )
        station_status = (
            data["station_status"] if data.get("station_status") else []
        )
        candidate_status = (
            data["candidate_status"] if data.get("candidate_status") else []
        )
        sub_con_codes = (
            data["sub_con_codes"] if data.get("sub_con_codes") else []
        )
        percentage_processed = (
            data["percentage_processed"]
            if data.get("percentage_processed")
            else 0
        )
        stations_processed_percentage = min(int(percentage_processed), 100)
        query_args = {}
        qs = qs.annotate(station_ids=station_id_query)

        stations_qs = Station.objects.filter(
            tally__id=tally_id,
            center__resultform__isnull=False,
        )
        if station_status:
            if len(station_status) == 1:
                station_status = station_status[0]
                if station_status == "active":
                    active = True
                else:
                    active = False
                if selected_station_ids:
                    stations_qs = stations_qs.filter(
                        id__in=selected_station_ids, active=active
                    )
                elif selected_center_ids:
                    stations_qs = stations_qs.filter(
                        center__id__in=selected_center_ids, active=active
                    )
                stations_qs = stations_qs.filter(active=active)
                selected_station_ids = (
                    [item.get("id") for item in stations_qs.values("id")]
                    if stations_qs.values("id")
                    else [0]
                )

        if stations_processed_percentage:
            if selected_station_ids:
                stations_qs = stations_qs.filter(
                    id__in=selected_station_ids,
                )
            elif selected_center_ids:
                stations_qs = stations_qs.filter(
                    center__id__in=selected_center_ids,
                )

            stations_qs = (
                stations_qs.values("id")
                .annotate(
                    total_result_forms=station_total_result_forms_sub_query,
                    total_result_forms_archived=station_total_result_forms_archived_sub_query,
                    processed_percentage=Round(
                        100
                        * F("total_result_forms_archived")
                        / F("total_result_forms"),
                        digits=2,
                    ),
                )
                .filter(
                    processed_percentage__gte=stations_processed_percentage
                )
            )
            selected_station_ids = (
                [item.get("id") for item in stations_qs]
                if stations_qs
                else [0]
            )

        if sub_race_type_names:
            sub_race_type_field = (
                "result_form__ballot__electrol_race__ballot_name__in"
            )
            query_args[sub_race_type_field] = sub_race_type_names

        if sub_con_codes:
            sub_con_code_field = (
                "result_form__center__sub_constituency__code__in"
            )
            query_args[sub_con_code_field] = sub_con_codes

        if election_level_names:
            election_level_field = (
                "result_form__ballot__electrol_race__election_level__in"
            )
            query_args[election_level_field] = election_level_names

        if ballot_status:
            if len(ballot_status) == 1:
                ballot_status = ballot_status[0]
                if ballot_status == "available_for_release":
                    available_for_release = True
                else:
                    available_for_release = False
                query_args["result_form__ballot__available_for_release"] = (
                    available_for_release
                )

        if candidate_status:
            if len(candidate_status) == 1:
                candidate_status = candidate_status[0]
                if candidate_status == "active":
                    active = True
                else:
                    active = False
                query_args["candidate__active"] = active

        qs = qs.filter(**query_args)
        if selected_station_ids or stations_processed_percentage:
            qs = qs.filter(Q(station_ids__in=selected_station_ids))

        elif selected_center_ids:
            qs = qs.filter(Q(result_form__center__id__in=selected_center_ids))

        qs = (
            qs.filter(candidate__full_name__isnull=False)
            .values("candidate_id")
            .annotate(
                candidate_number=F("candidate__candidate_id"),
                candidate_name=F("candidate__full_name"),
                total_votes=Sum("votes"),
                gender=F("result_form__gender"),
                center_code=F("result_form__center__code"),
                center_name=F("result_form__center__name"),
                office_number=F("result_form__office__number"),
                office_name=F("result_form__office__name"),
                station_number=F("result_form__station_number"),
                electrol_race_id=F("result_form__ballot__electrol_race__id"),
                election_level=F(
                    "result_form__ballot__electrol_race__election_level"
                ),
                sub_race_type=F(
                    "result_form__ballot__electrol_race__ballot_name"
                ),
                sub_con_name=F("result_form__center__sub_constituency__name"),
                sub_con_code=F("result_form__center__sub_constituency__code"),
                order=F("candidate__order"),
                ballot_number=F("candidate__ballot__number"),
                candidate_status=Case(
                    When(candidate__active=True, then=V("enabled")),
                    default=V("disabled"),
                    output_field=CharField(),
                ),
            )
        )

    else:
        qs = (
            qs.filter(candidate__full_name__isnull=False)
            .values("candidate_id")
            .annotate(
                candidate_number=F("candidate__candidate_id"),
                candidate_name=F("candidate__full_name"),
                total_votes=Sum("votes"),
                gender=F("result_form__gender"),
                center_code=F("result_form__center__code"),
                center_name=F("result_form__center__name"),
                office_number=F("result_form__office__number"),
                office_name=F("result_form__office__name"),
                station_number=F("result_form__station_number"),
                electrol_race_id=F("result_form__ballot__electrol_race__id"),
                election_level=F(
                    "result_form__ballot__electrol_race__election_level"
                ),
                sub_race_type=F(
                    "result_form__ballot__electrol_race__ballot_name"
                ),
                sub_con_name=F("result_form__center__sub_constituency__name"),
                sub_con_code=F("result_form__center__sub_constituency__code"),
                order=F("candidate__order"),
                ballot_number=F("candidate__ballot__number"),
                candidate_status=Case(
                    When(candidate__active=True, then=V("enabled")),
                    default=V("disabled"),
                    output_field=CharField(),
                ),
            )
        )

    return qs


def export_results_queryset(tally_id, qs, data=None):
    """
    Genarate votes per candidate results for export.

    :param tally_id: The tally id.
    :param qs: The result parent queryset.
    :param data: An array of dicts containing centers and stations
        id's to filter out from the queryset.

    returns: The votes per candidate queryset.
    """
    station_id_query = Subquery(
        Station.objects.filter(
            tally__id=tally_id,
            center__code=OuterRef("result_form__center__code"),
            station_number=OuterRef("result_form__station_number"),
        ).values("id")[:1],
        output_field=IntegerField(),
    )

    station_total_result_forms_sub_query = Subquery(
        ResultForm.objects.filter(
            tally__id=tally_id,
            center__code=OuterRef("center__code"),
            station_number=OuterRef("station_number"),
        )
        .values("center__code", "station_number")
        .annotate(
            total_result_forms=Count("barcode", distinct=True),
        )
        .values("total_result_forms")[:1],
        output_field=IntegerField(),
    )

    station_total_result_forms_archived_sub_query = Subquery(
        ResultForm.objects.filter(
            tally__id=tally_id,
            center__code=OuterRef("center__code"),
            station_number=OuterRef("station_number"),
        )
        .values("center__code", "station_number")
        .annotate(
            total_result_forms_archived=Count(
                "barcode",
                distinct=True,
                filter=Q(form_state=FormState.ARCHIVED),
            ),
        )
        .values("total_result_forms_archived")[:1],
        output_field=IntegerField(),
    )

    if data:
        selected_center_ids = (
            data["select_1_ids"] if data.get("select_1_ids") else []
        )
        selected_station_ids = (
            data["select_2_ids"] if data.get("select_2_ids") else []
        )
        election_level_names = (
            data["election_level_names"]
            if data.get("election_level_names")
            else []
        )
        sub_race_type_names = (
            data["sub_race_type_names"]
            if data.get("sub_race_type_names")
            else []
        )
        ballot_status = (
            data["ballot_status"] if data.get("ballot_status") else []
        )
        station_status = (
            data["station_status"] if data.get("station_status") else []
        )
        candidate_status = (
            data["candidate_status"] if data.get("candidate_status") else []
        )
        sub_con_codes = (
            data["sub_con_codes"] if data.get("sub_con_codes") else []
        )
        percentage_processed = (
            data["percentage_processed"]
            if data.get("percentage_processed")
            else 0
        )
        stations_processed_percentage = min(int(percentage_processed), 100)
        query_args = {}
        qs = qs.annotate(station_ids=station_id_query)

        stations_qs = Station.objects.filter(
            tally__id=tally_id,
            center__resultform__isnull=False,
        )
        if station_status:
            if len(station_status) == 1:
                station_status = station_status[0]
                if station_status == "active":
                    active = True
                else:
                    active = False
                if selected_station_ids:
                    stations_qs = stations_qs.filter(
                        id__in=selected_station_ids, active=active
                    )
                elif selected_center_ids:
                    stations_qs = stations_qs.filter(
                        center__id__in=selected_center_ids, active=active
                    )
                stations_qs = stations_qs.filter(active=active)
                selected_station_ids = (
                    [item.get("id") for item in stations_qs.values("id")]
                    if stations_qs.values("id")
                    else [0]
                )

        if stations_processed_percentage:
            if selected_station_ids:
                stations_qs = stations_qs.filter(
                    id__in=selected_station_ids,
                )
            elif selected_center_ids:
                stations_qs = stations_qs.filter(
                    center__id__in=selected_center_ids,
                )

            stations_qs = (
                stations_qs.values("id")
                .annotate(
                    total_result_forms=station_total_result_forms_sub_query,
                    total_result_forms_archived=station_total_result_forms_archived_sub_query,
                    processed_percentage=Round(
                        100
                        * F("total_result_forms_archived")
                        / F("total_result_forms"),
                        digits=2,
                    ),
                )
                .filter(
                    processed_percentage__gte=stations_processed_percentage
                )
            )
            selected_station_ids = (
                [item.get("id") for item in stations_qs]
                if stations_qs
                else [0]
            )

        if sub_race_type_names:
            sub_race_type_field = (
                "result_form__ballot__electrol_race__ballot_name__in"
            )
            query_args[sub_race_type_field] = sub_race_type_names

        if election_level_names:
            election_level_field = (
                "result_form__ballot__electrol_race__election_level__in"
            )
            query_args[election_level_field] = election_level_names

        if ballot_status:
            if len(ballot_status) == 1:
                ballot_status = ballot_status[0]
                if ballot_status == "available_for_release":
                    available_for_release = True
                else:
                    available_for_release = False
                query_args["result_form__ballot__available_for_release"] = (
                    available_for_release
                )

        if candidate_status:
            if len(candidate_status) == 1:
                candidate_status = candidate_status[0]
                if candidate_status == "active":
                    active = True
                else:
                    active = False
                query_args["candidate__active"] = active

        if sub_con_codes:
            sub_con_code_field = (
                "result_form__center__sub_constituency__code__in"
            )
            query_args[sub_con_code_field] = sub_con_codes

        qs = qs.filter(**query_args)
        if selected_station_ids or stations_processed_percentage:
            qs = qs.filter(Q(station_ids__in=selected_station_ids))

        elif selected_center_ids:
            qs = qs.filter(Q(result_form__center__id__in=selected_center_ids))

        qs = (
            qs.filter(candidate__full_name__isnull=False)
            .values("candidate_id")
            .annotate(
                candidate_name=F("candidate__full_name"),
                total_votes=Sum("votes"),
                electrol_race_id=F("result_form__ballot__electrol_race__id"),
                election_level=F(
                    "result_form__ballot__electrol_race__election_level"
                ),
                sub_race_type=F(
                    "result_form__ballot__electrol_race__ballot_name"
                ),
                sub_con_name=F("result_form__center__sub_constituency__name"),
            )
        )

    else:
        qs = (
            qs.filter(candidate__full_name__isnull=False)
            .values("candidate_id")
            .annotate(
                candidate_name=F("candidate__full_name"),
                total_votes=Sum("votes"),
                electrol_race_id=F("result_form__ballot__electrol_race__id"),
                election_level=F(
                    "result_form__ballot__electrol_race__election_level"
                ),
                sub_race_type=F(
                    "result_form__ballot__electrol_race__ballot_name"
                ),
                sub_con_name=F("result_form__center__sub_constituency__name"),
            )
        )

    return qs


def duplicate_results_queryset(tally_id, qs, data=None):
    """
    Genarate a report of duplicate results per result form.

    :param tally_id: The tally id.
    :param qs: The result form parent queryset.
    :param data: An array of dicts containing centers and stations
        id's to filter out from the queryset.

    returns: The duplicate results queryset.
    """
    station_id_query = Subquery(
        Station.objects.filter(
            tally__id=tally_id,
            center__code=OuterRef("center__code"),
            station_number=OuterRef("station_number"),
        ).values("id")[:1],
        output_field=IntegerField(),
    )

    if data:
        selected_center_ids = (
            data["select_1_ids"] if len(data["select_1_ids"]) else [0]
        )
        selected_station_ids = (
            data["select_2_ids"] if len(data["select_2_ids"]) else [0]
        )

        qs = qs.annotate(station_ids=station_id_query).filter(
            ~Q(center__id__in=selected_center_ids)
            & ~Q(station_ids__in=selected_station_ids)
        )

        result_form_votes_registrants_query = Subquery(
            Result.objects.filter(
                result_form__tally__id=tally_id,
                result_form__form_state=FormState.ARCHIVED,
                entry_version=EntryVersion.FINAL,
                active=True,
                result_form__center__code=OuterRef("center__code"),
                result_form__station_number=OuterRef("station_number"),
            )
            .values("result_form__barcode")
            .annotate(total_votes=Coalesce(Sum("votes"), V(0)))
            .values("total_votes")[:1],
            output_field=IntegerField(),
        )

        qs = (
            qs.values("barcode")
            .annotate(
                ballot_number=F("ballot__number"),
                center_code=F("center__code"),
                state=F("form_state"),
                station_number=F("station_number"),
                station_id=station_id_query,
                votes=result_form_votes_registrants_query,
                office=F('center__office__name'),
            )
            .distinct()
        )
    else:
        result_form_votes_registrants_query = Subquery(
            Result.objects.filter(
                result_form__tally__id=tally_id,
                result_form__form_state=FormState.ARCHIVED,
                entry_version=EntryVersion.FINAL,
                active=True,
                result_form__center__code=OuterRef("center__code"),
                result_form__station_number=OuterRef("station_number"),
            )
            .values("result_form__barcode")
            .annotate(total_votes=Coalesce(Sum("votes"), V(0)))
            .values("total_votes")[:1],
            output_field=IntegerField(),
        )

        qs = (
            qs.values("barcode")
            .annotate(
                ballot_number=F("ballot__number"),
                center_code=F("center__code"),
                state=F("form_state"),
                station_id=station_id_query,
                station_number=F("station_number"),
                votes=result_form_votes_registrants_query,
                office=F('center__office__name'),
            )
            .distinct()
        )

    return qs


def filter_candidates_votes_queryset(qs, data=None):
    """
    Filter candidates votes report.

    :param qs: The votes queryset.
    :param data: An array of dicts containing centers and stations
        id's to filter out from the queryset.

    returns: The queryset containing candidates votes grouped by
        candidate name.
    """
    if data:
        selected_center_ids = (
            data["select_1_ids"] if len(data["select_1_ids"]) else [0]
        )
        # TODO: Find a way to filter by station ids
        # selected_station_ids =\
        #     data['select_2_ids'] if len(data['select_2_ids']) else [0]
        qs = qs.filter(~Q(center_ids__contains=selected_center_ids))

    return qs


def generate_progressive_report(
    tally_id, report_column_name, region_id=None, constituency_id=None
):
    """
    Genarate progressive report of candidates by votes.

    :param tally_id: The result form tally.
    :param report_column_name: The result form report column name.
    :param region_id: The result form region id.
    :param constituency_id: The result form constituency id.

    returns: The candidates votes stats based on an administrative area.
    """
    qs = Result.objects.filter(
        result_form__tally__id=tally_id,
        result_form__form_state=FormState.ARCHIVED,
        entry_version=EntryVersion.FINAL,
        active=True,
    )

    if region_id:
        qs = qs.filter(result_form__office__region__id=region_id)

    if constituency_id:
        qs = qs.filter(result_form__center__constituency__id=constituency_id)

    qs = (
        qs.annotate(
            name=F(report_column_name),
            admin_area_id=F("result_form__office__region__id"),
        )
        .values(
            "name",
            "admin_area_id",
        )
        .annotate(
            total_candidates=Count("candidate__id", distinct=True),
            total_votes=Sum("votes"),
        )
    )

    if region_id:
        qs = qs.annotate(
            constituency_id=F("result_form__center__constituency__id"),
            sub_constituency_id=F("result_form__center__sub_constituency__id"),
        )

    return qs


def generate_progressive_report_queryset(qs, data=None, **kwargs):
    """
    Genarate progressive report of candidates by votes.

    :param tally_id: The result form tally.
    :param report_column_name: The result form report column name.
    :param region_id: The result form region id.
    :param constituency_id: The result form constituency id.

    returns: The candidates votes stats based on an administrative area.
    """
    column_name = "result_form__center__office__region__name"
    column_id = "result_form__center__office__region__id"
    admin_area_id = kwargs.get("region_id")
    constituency_id = kwargs.get("constituency_id")

    if admin_area_id and constituency_id:
        column_name = "result_form__center__sub_constituency__code"
        column_id = "result_form__center__sub_constituency__id"
    elif admin_area_id:
        column_id = "result_form__center__constituency__id"
        column_name = "result_form__center__constituency__name"

    if admin_area_id:
        qs = qs.filter(result_form__office__region__id=admin_area_id)

    if constituency_id:
        qs = qs.filter(result_form__center__constituency__id=constituency_id)

    if data and len(
        [
            data
            for d in data
            if len(d["select_1_ids"]) or len(d["select_2_ids"])
        ]
    ):
        parent_qs = qs
        for item in data:
            region_id = item["region_id"]
            constituency_ids = (
                item["select_1_ids"] if len(item["select_1_ids"]) else [0]
            )
            sub_constituency_ids = (
                item["select_2_ids"] if len(item["select_2_ids"]) else [0]
            )
            con_ids = constituency_ids
            filter_by_constituency_id = ~Q(
                result_form__center__constituency__id__in=con_ids
            )
            subcn_ids = sub_constituency_ids
            filter_by_sub_constituency_id = ~Q(
                result_form__center__sub_constituency__id__in=subcn_ids
            )

            if admin_area_id and constituency_id:
                current_qs = parent_qs.filter(
                    result_form__center__sub_constituency__id=region_id
                )
            elif admin_area_id:
                current_qs = parent_qs.filter(
                    result_form__center__constituency__id=region_id
                )
            else:
                current_qs = parent_qs.filter(
                    result_form__center__office__region__id=region_id
                )

            current_qs = (
                current_qs.annotate(
                    admin_area_name=F(column_name), admin_area_id=F(column_id)
                )
                .values(
                    "admin_area_name",
                    "admin_area_id",
                )
                .annotate(
                    total_candidates=Count(
                        "candidate__id",
                        distinct=True,
                        filter=(
                            filter_by_constituency_id
                            & filter_by_sub_constituency_id
                        ),
                    ),
                    total_votes=Coalesce(
                        Sum(
                            "votes",
                            filter=(
                                filter_by_constituency_id
                                & filter_by_sub_constituency_id
                            ),
                        ),
                        V(0),
                    ),
                    region_id=F("result_form__center__office__region__id"),
                    constituencies_ids=ArrayAgg(
                        "result_form__center__constituency__id", distinct=True
                    ),
                    sub_constituencies_ids=ArrayAgg(
                        "result_form__center__sub_constituency__id",
                        distinct=True,
                    ),
                )
            )

            if admin_area_id:
                current_qs = current_qs.annotate(
                    constituency_id=F("result_form__center__constituency__id"),
                    sub_constituency_id=F(
                        "result_form__center__sub_constituency__id"
                    ),
                )
            qs = (
                qs.union(current_qs)
                if not isinstance(qs[0], Result)
                else current_qs
            )
    else:
        qs = (
            qs.annotate(
                admin_area_name=F(column_name), admin_area_id=F(column_id)
            )
            .values(
                "admin_area_name",
                "admin_area_id",
            )
            .annotate(
                total_candidates=Count("candidate__id", distinct=True),
                total_votes=Sum("votes"),
                region_id=F("result_form__center__office__region__id"),
                constituencies_ids=ArrayAgg(
                    "result_form__center__constituency__id", distinct=True
                ),
                sub_constituencies_ids=ArrayAgg(
                    "result_form__center__sub_constituency__id", distinct=True
                ),
            )
        )

        if admin_area_id:
            qs = qs.annotate(
                constituency_id=F("result_form__center__constituency__id"),
                sub_constituency_id=F(
                    "result_form__center__sub_constituency__id"
                ),
            )

    return qs


def get_admin_areas_with_forms_in_audit(
    tally_id, report_column_name, region_id=None, constituency_id=None
):
    """
    Genarate a report of stations and centers with result forms in audit state.

    :param tally_id: The reconciliation forms tally.
    :param report_column_name: The result form report column name.
    :param region_id: The result form report region id used for filtering.
    :param constituency_id: The result form report constituency id
        used for filtering.

    returns: The stations and centers report grouped by the report column name.
    """
    qs = ResultForm.objects.filter(
        tally__id=tally_id, form_state=FormState.AUDIT
    )
    if region_id:
        qs = qs.filter(office__region__id=region_id)
    if constituency_id:
        qs = qs.filter(center__constituency__id=constituency_id)
    qs = (
        qs.annotate(
            admin_area_name=F(report_column_name),
            region_id=F("office__region__id"),
        )
        .values(
            "admin_area_name",
            "region_id",
        )
        .annotate(
            number_of_centers_in_audit_state=Count("center"),
            number_of_stations_in_audit_state=Count("station_number"),
            total_num_of_centers_and_stations_in_audit=ExpressionWrapper(
                F("number_of_centers_in_audit_state")
                + F("number_of_stations_in_audit_state"),
                output_field=IntegerField(),
            ),
        )
    )

    if region_id:
        qs = qs.annotate(
            constituency_id=F("center__constituency__id"),
            sub_constituency_id=F("center__sub_constituency__id"),
        )

    return qs


def custom_queryset_filter(tally_id, qs, data=None, **kwargs):
    """
    Filter queryset by tally_id, region_id and constituency_ids.

    :param tally_id: The reconciliation forms tally.
    :param region_id: The region id for filtering the reconciliation forms.
    :param constituency_ids: The constituency ids to exclude from the queryset.

    returns: The filtered queryset.
    """
    admin_area_id = kwargs.get("region_id")
    constituency_id = kwargs.get("constituency_id")
    report_type = kwargs.get("report_type")
    column_name = "result_form__office__region__name"
    column_id = "result_form__center__office__region__id"
    turnout_report_type = report_types[1]
    summary_report_type = report_types[2]

    if admin_area_id and constituency_id:
        column_name = "result_form__center__sub_constituency__code"
        column_id = "result_form__center__sub_constituency__id"
    elif admin_area_id:
        column_name = "result_form__center__constituency__name"
        column_id = "result_form__center__constituency__id"

    if admin_area_id:
        qs = qs.filter(result_form__office__region__id=admin_area_id)

    if constituency_id:
        qs = qs.filter(result_form__center__constituency__id=constituency_id)

    qs = qs.filter(result_form__tally__id=tally_id,
                   result_form__form_state=FormState.ARCHIVED,
                   entry_version=EntryVersion.FINAL)

    if data:
        parent_qs = qs
        for item in data:
            item["region_id"]
            constituency_ids = (
                item["select_1_ids"] if len(item["select_1_ids"]) else [0]
            )
            sub_constituency_ids = (
                item["select_2_ids"] if len(item["select_2_ids"]) else [0]
            )

            con_ids = constituency_ids
            filter_by_constituency_id = ~Q(
                result_form__center__constituency__id__in=con_ids
            )
            subcn_ids = sub_constituency_ids
            filter_by_sub_constituency_id = ~Q(
                result_form__center__sub_constituency__id__in=subcn_ids
            )

            # This seems to be already filtered out and the target
            # seems not to be correct.
            # if admin_area_id and constituency_id:
            #     current_qs = parent_qs.get_registrants_and_votes_type().
            #        filter(
            #         result_form__center__sub_constituency__id=region_id
            #     )
            # elif admin_area_id:
            #     current_qs = parent_qs.get_registrants_and_votes_type().
            #     filter(
            #         result_form__center__constituency__id=region_id
            #     )
            # else:
            #     current_qs = parent_qs.get_registrants_and_votes_type()
            #        .filter(
            #         result_form__center__office__region__id=region_id
            #     )

            current_qs = (
                parent_qs.get_registrants_and_votes_type().filter()
                .annotate(name=F(column_name), admin_area_id=F(column_id))
                .values(
                    "name",
                    "admin_area_id",
                )
                .annotate(
                    region_id=F("result_form__office__region__id"),
                    constituency_id=F("result_form__center__constituency__id"),
                )
            )

            if report_type == turnout_report_type:
                current_qs = (
                    current_qs.annotate(
                        number_of_voters_voted=Coalesce(
                            Sum(
                                "number_valid_votes",
                                filter=(
                                    filter_by_constituency_id
                                    & filter_by_sub_constituency_id
                                ),
                                default=V(0),
                            ),
                            V(0),
                        )
                    )
                    .annotate(
                        total_number_of_registrants=Sum(
                            "result_form__center__stations__registrants",
                            filter=(
                                filter_by_constituency_id
                                & filter_by_sub_constituency_id
                            ),
                            default=V(0),
                        )
                    )
                    .annotate(
                        total_number_of_ballots_used=Coalesce(
                            Sum(
                                ExpressionWrapper(
                                    F("number_valid_votes")
                                    + F("number_invalid_votes"),
                                    output_field=IntegerField(),
                                ),
                                filter=(
                                    filter_by_constituency_id
                                    & filter_by_sub_constituency_id
                                ),
                            ),
                            V(0),
                        )
                    )
                    .annotate(
                        turnout_percentage=Coalesce(
                            ExpressionWrapper(
                                V(100)
                                * F("total_number_of_ballots_used")
                                / F("total_number_of_registrants"),
                                output_field=IntegerField(),
                            ),
                            V(0),
                        )
                    )
                    .annotate(
                        male_voters=Coalesce(
                            Sum(
                                "number_valid_votes",
                                filter=(
                                    Q(voters_gender_type=0)
                                    & filter_by_constituency_id
                                    & filter_by_sub_constituency_id
                                ),
                            ),
                            V(0),
                        )
                    )
                    .annotate(
                        female_voters=Coalesce(
                            Sum(
                                "number_valid_votes",
                                filter=(
                                    Q(voters_gender_type=1)
                                    & filter_by_constituency_id
                                    & filter_by_sub_constituency_id
                                ),
                            ),
                            V(0),
                        )
                    )
                    .annotate(
                        constituencies_ids=ArrayAgg(
                            "result_form__center__constituency__id",
                            distinct=True,
                        )
                    )
                    .annotate(
                        sub_constituencies_ids=ArrayAgg(
                            "result_form__center__sub_constituency__id",
                            distinct=True,
                        )
                    )
                )

            if report_type == summary_report_type:
                current_qs = current_qs.annotate(
                    number_valid_votes=Coalesce(
                        Sum(
                            "number_valid_votes",
                            filter=(
                                filter_by_constituency_id
                                & filter_by_sub_constituency_id
                            ),
                            default=V(0),
                        ),
                        V(0),
                    ),
                    number_invalid_votes=Coalesce(
                        Sum(
                            "number_invalid_votes",
                            filter=(
                                filter_by_constituency_id
                                & filter_by_sub_constituency_id
                            ),
                            default=V(0),
                        ),
                        V(0),
                    ),
                    constituencies_ids=ArrayAgg(
                        "result_form__center__constituency__id", distinct=True
                    ),
                    sub_constituencies_ids=ArrayAgg(
                        "result_form__center__sub_constituency__id",
                        distinct=True,
                    ),
                )

            qs = (
                qs.union(current_qs)
                if not isinstance(qs[0], ReconciliationForm)
                else current_qs
            )
    else:
        qs = (
            qs.get_registrants_and_votes_type()
            .filter()
            .annotate(name=F(column_name), admin_area_id=F(column_id))
            .values(
                "name",
                "admin_area_id",
            )
            .annotate(
                region_id=F("result_form__office__region__id"),
                constituency_id=F("result_form__center__constituency__id"),
            )
        )

        if report_type == turnout_report_type:
            qs = (
                qs.annotate(number_of_voters_voted=Sum("number_valid_votes"))
                .annotate(
                    total_number_of_registrants=Sum(
                        "result_form__center__stations__registrants"
                    )
                )
                .annotate(
                    total_number_of_ballots_used=Sum(
                        ExpressionWrapper(
                            F("number_valid_votes")
                            + F("number_invalid_votes"),
                            output_field=IntegerField(),
                        )
                    )
                )
                .annotate(
                    turnout_percentage=ExpressionWrapper(
                        V(100)
                        * F("total_number_of_ballots_used")
                        / F("total_number_of_registrants"),
                        output_field=IntegerField(),
                    )
                )
                .annotate(
                    male_voters=Coalesce(
                        Sum(
                            "number_valid_votes",
                            filter=Q(voters_gender_type=0),
                        ),
                        V(0),
                    )
                )
                .annotate(
                    female_voters=Coalesce(
                        Sum(
                            "number_valid_votes",
                            filter=Q(voters_gender_type=1),
                        ),
                        V(0),
                    )
                )
                .annotate(
                    constituencies_ids=ArrayAgg(
                        "result_form__center__constituency__id", distinct=True
                    )
                )
                .annotate(
                    sub_constituencies_ids=ArrayAgg(
                        "result_form__center__sub_constituency__id",
                        distinct=True,
                    )
                )
            )

        if report_type == summary_report_type:
            qs = qs.annotate(
                number_valid_votes=Sum("number_valid_votes"),
                number_invalid_votes=Sum("number_invalid_votes"),
                constituencies_ids=ArrayAgg(
                    "result_form__center__constituency__id", distinct=True
                ),
                sub_constituencies_ids=ArrayAgg(
                    "result_form__center__sub_constituency__id", distinct=True
                ),
            )

    return qs


def build_select_options(qs, ids=None):
    if ids is None:
        ids = []
    select = 'selected="selected"'

    return [
        str(
            "<option "
            f"{select if str(item[0]) in ids else ''}"
            " value="
            f"{item[0]}"
            ">"
            f"{item[1]}"
            "</option>"
        )
        for item in list(qs)
    ]


def get_centers_stations(request):
    """
    Retrieves stations that belong to the centers available in the request.

    :param request: The request object containing the list of center ids.

    returns: A JSON response of the centers stations ids
    """
    data = ast.literal_eval(request.GET.get("data"))
    center_ids = data.get("center_ids")
    tally_id = data.get("tally_id")

    return JsonResponse(
        {
            "station_ids": list(
                Station.objects.filter(
                    tally__id=tally_id, center__id__in=center_ids
                ).values_list("id", flat=True)
            )
        }
    )


def get_export(request):
    """
    Generates and returns a PowerPoint export based on the filter
    values provided
    """
    data = ast.literal_eval(request.GET.get("data"))
    tally_id = data.get("tally_id")
    limit = parse_int(data.get("export_number"))
    export_type = data.get("exportType")
    center_ids = data.get("select_1_ids")
    station_ids = data.get("select_2_ids")
    election_level_names = data.get("election_level_names")
    sub_race_type_names = data.get("sub_race_type_names")
    ballot_status = data.get("ballot_status")
    station_status = data.get("station_status")
    candidate_status = data.get("candidate_status")
    percentage_processed = data.get("percentage_processed")
    sub_con_codes = data.get("sub_con_codes")
    filters_applied = (
        center_ids
        or station_ids
        or election_level_names
        or sub_race_type_names
        or ballot_status
        or station_status
        or candidate_status
        or percentage_processed
        or sub_con_codes
    )

    qs = Result.objects.select_related(
        "candidate",
    )
    qs = qs.filter(
        result_form__tally__id=tally_id,
        result_form__form_state=FormState.ARCHIVED,
        entry_version=EntryVersion.FINAL,
        active=True,
    )
    qs = get_filtered_candidate_votes(
        tally_id, qs, data=data if filters_applied else None
    )

    if export_type == "PPT" and qs.count() != 0:
        result = create_ppt_export(qs, tally_id=tally_id, limit=limit)
        return result
    return HttpResponse("Not found")


def electrol_race_has_results(results_qs, electrol_race):
    results_qs = results_qs.filter(
        election_level=electrol_race.election_level,
        sub_race_type=electrol_race.ballot_name,
    ).order_by("-votes")
    return results_qs.count() != 0


def create_ppt_export(qs, tally_id=None, limit=None):
    electrol_race_ids = list(
        set([result.get("electrol_race_id") for result in qs])
    )
    filtered_electrol_races = ElectrolRace.objects.filter(
        tally__id=tally_id, id__in=electrol_race_ids
    )

    headers = create_results_power_point_headers(
        tally_id, filtered_electrol_races, qs
    )

    powerpoint_data = []
    race_bg_img_root_path = settings.MEDIA_ROOT

    races_ids_to_total_votes_mapping = {}
    for electral_race_id in filtered_electrol_races.values_list(
        "id", flat=True
    ):
        races_ids_to_total_votes_mapping[electral_race_id] = sum(
            [
                q.get("total_votes")
                for q in qs
                if q.get("electrol_race_id") == electral_race_id
            ]
        )

    for electrol_race in filtered_electrol_races:
        if electrol_race_has_results(qs, electrol_race):
            data = [
                r for r in qs if r.get("electrol_race_id") == electrol_race.id
            ]
            sorted_results = sorted(data, key=lambda x: -x["total_votes"])
            body_data = sorted_results
            if limit:
                body_data = body_data[:limit]
            body = [
                {
                    "candidate_name": item.get("candidate_name"),
                    "total_votes": item.get("total_votes"),
                    "valid_votes": races_ids_to_total_votes_mapping.get(
                        item.get("electrol_race_id")
                    ),
                }
                for item in body_data
            ]
            header = headers.get(
                f"{electrol_race.election_level}_{electrol_race.ballot_name}"
            )
            file_path = (
                f"{race_bg_img_root_path}/{electrol_race.background_image}"
            )
            powerpoint_data.append(
                {
                    "header": header,
                    "body": body,
                    "background_image": (
                        file_path if os.path.isfile(file_path) else None
                    ),
                }
            )

    # Create a new PowerPoint presentation
    prs = Presentation()

    # Set the cover page
    create_results_power_point_cover_page(prs)

    for data in powerpoint_data:
        # Create the summary slide
        create_results_power_point_summary_slide(
            prs, power_point_race_data=data
        )

        # Create the candidates slides
        create_results_power_point_candidates_results_slide(
            prs, power_point_race_data=data, limit=limit
        )

    # Save the presentation to a file
    formatted_datestring = datetime.date.today().strftime("%Y%m%d")
    response = save_ppt_presentation_to_file(
        prs, f"election_results_{formatted_datestring}.pptx"
    )

    return response


def save_ppt_presentation_to_file(prs, file_name):
    pptx_stream = BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)

    # Create an HTTP response with the PowerPoint file
    content_type = str(
        "application/vnd.openxmlformats-officedocument"
        ".presentationml.presentation"
    )
    response = HttpResponse(pptx_stream.getvalue(), content_type=content_type)
    response["Content-Disposition"] = f"attachment; filename={file_name}"

    return response


def create_results_power_point_cover_page(prs):
    # Set the cover page
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    background_image = settings.CANDIDATE_RESULTS_PPT_COVER_PAGE_BCK_IMG_PATH

    # Set background image if provided
    if background_image:
        slide.shapes.add_picture(
            background_image,
            Inches(0),
            Inches(0),
            prs.slide_width,
            prs.slide_height,
        )

    # Add date of creation
    date_of_creation = date.today().strftime("%B %d, %Y")
    date_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(8), Inches(0.5)
    ).text_frame
    date_shape.text = "Date: " + date_of_creation
    date_shape.paragraphs[0].alignment = PP_ALIGN.CENTER
    date_shape.paragraphs[0].runs[0].font.bold = True

    return


def create_results_power_point_summary_slide(prs, power_point_race_data):
    # Create the summary slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    background_image = power_point_race_data["background_image"]
    election_level_name = power_point_race_data["header"][
        "election_level"
    ].capitalize()
    sub_race = power_point_race_data["header"]["sub_race_type"].capitalize()

    summary_slide_title = (
        f"{election_level_name} {sub_race} Election Summary Results"
    )
    # Set background image if provided
    if background_image:
        slide.shapes.add_picture(
            background_image,
            Inches(0),
            Inches(0),
            prs.slide_width,
            prs.slide_height,
        )

    # Access the title placeholder and set its text
    if slide.shapes.title:
        slide.shapes.title.text = summary_slide_title
        # Apply formatting to the title
        title_frame = slide.shapes.title.text_frame
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    else:
        # Add title text box for the summary slide
        title_text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.7), prs.slide_width - Inches(1), Inches(0.5)
        )
        title_text_frame = title_text_box.text_frame

        # Set title properties
        title_text_frame.text = summary_slide_title
        title_text_frame.word_wrap = True
        title_text_frame.margin_left = 0
        title_text_frame.margin_right = 0
        title_text_frame.margin_top = 0
        title_text_frame.margin_bottom = 0

        # Set title font properties
        title_text_frame.clear()  # Clear existing paragraphs
        p = title_text_frame.add_paragraph()
        p.text = summary_slide_title
        p.font.bold = True
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

    # Create a table shape for the summary data
    summary_table = slide.shapes.add_table(
        rows=9,
        cols=2,
        left=Inches(0.5),
        top=Inches(1.7),
        width=Inches(9),
        height=Inches(2),
    ).table

    # Set the column widths for the summary table
    column_widths = [Inches(4.5), Inches(4.5)]
    for i, width in enumerate(column_widths):
        summary_table.columns[i].width = width

    # Populate the summary table with data
    summary_table.cell(0, 0).text = "Election Level"
    summary_table.cell(0, 1).text = str(
        power_point_race_data["header"]["election_level"]
    )
    summary_table.cell(1, 0).text = "Name of Race"
    summary_table.cell(1, 1).text = str(
        power_point_race_data["header"]["sub_race_type"]
    )
    summary_table.cell(2, 0).text = "Stations Expected"
    summary_table.cell(
        2, 1
    ).text = f"{power_point_race_data['header']['stations_expected']:,.0f}"
    summary_table.cell(3, 0).text = "Stations Processed"
    summary_table.cell(3, 1).text = str(
        power_point_race_data["header"]["stations_processed"]
    )
    summary_table.cell(4, 0).text = "Percentage of Stations Processed"
    percentage_of_stations_processed = power_point_race_data["header"][
        "percentage_of_stations_processed"
    ]
    summary_table.cell(4, 1).text = f"{percentage_of_stations_processed}%"
    summary_table.cell(5, 0).text = "Results Status"
    summary_table.cell(5, 1).text = power_point_race_data["header"][
        "results_status"
    ]
    summary_table.cell(6, 0).text = "Registrants"
    registrants_in_processed_stations = power_point_race_data["header"][
        "registrants_in_processed_stations"
    ]
    summary_table.cell(6, 1).text = f"{registrants_in_processed_stations:,.0f}"
    summary_table.cell(7, 0).text = "Ballots Cast"
    summary_table.cell(7, 1).text = str(
        power_point_race_data["header"]["voters_in_counted_stations"]
    )
    summary_table.cell(8, 0).text = "Turnout"
    summary_table.cell(
        8, 1
    ).text = f"{power_point_race_data['header']['percentage_turnout']}%"

    return


def create_results_power_point_candidates_results_slide(
    prs, power_point_race_data, limit
):
    background_image = power_point_race_data["background_image"]
    candidates = power_point_race_data["body"]
    num_candidates = len(candidates)
    max_candidates_per_slide = 10
    num_slides = (num_candidates - 1) // max_candidates_per_slide + 1
    candidate_rank = 0
    election_level_name = power_point_race_data["header"][
        "election_level"
    ].capitalize()
    sub_race = power_point_race_data["header"]["sub_race_type"].capitalize()

    for slide_num in range(num_slides):
        # Create a new candidates slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        candidates_result_slide_title = (
            f"Showing all {election_level_name} {sub_race} Election Results"
        )
        if limit:
            candidates_result_slide_title = str(
                f"Top {limit} Leading {election_level_name} {sub_race} "
                "Election Results"
            )
        # Set background image if provided
        if background_image:
            slide.shapes.add_picture(
                background_image,
                Inches(0),
                Inches(0),
                prs.slide_width,
                prs.slide_height,
            )

        # Access the title placeholder and set its text
        if slide.shapes.title:
            slide.shapes.title.text = candidates_result_slide_title
            # Apply formatting to the title
            title_frame = slide.shapes.title.text_frame
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.size = Pt(24)
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        else:
            # Add title text box for the candidates slide
            title_text_box = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(0.7),
                prs.slide_width - Inches(1),
                Inches(0.5),
            )
            title_text_frame = title_text_box.text_frame

            # Set title properties
            title_text_frame.text = candidates_result_slide_title
            title_text_frame.word_wrap = True
            title_text_frame.margin_left = 0
            title_text_frame.margin_right = 0
            title_text_frame.margin_top = 0
            title_text_frame.margin_bottom = 0

            # Set title font properties
            title_text_frame.clear()  # Clear existing paragraphs
            p = title_text_frame.add_paragraph()
            p.text = candidates_result_slide_title
            p.font.bold = True
            p.font.size = Pt(24)
            p.alignment = PP_ALIGN.CENTER

        # Calculate the number of candidates for the current slide
        start_index = slide_num * max_candidates_per_slide
        end_index = start_index + max_candidates_per_slide
        candidates_slice = candidates[start_index:end_index]

        # Create a table shape for the candidates data
        candidates_table = slide.shapes.add_table(
            rows=len(candidates_slice) + 1,
            cols=5,
            left=Inches(0.5),
            top=Inches(1.7),
            width=Inches(9),
            height=Inches(2),
        ).table

        # Set the column widths for the candidates table
        column_widths = [
            Inches(1),
            Inches(3),
            Inches(2),
            Inches(1.5),
            Inches(1.5),
        ]
        for i, width in enumerate(column_widths):
            candidates_table.columns[i].width = width

        # Populate the candidates table with data
        candidates_table.cell(0, 0).text = "Rank"
        candidates_table.cell(0, 1).text = "Name"
        candidates_table.cell(0, 2).text = "Votes"
        candidates_table.cell(0, 3).text = "Total Votes"
        candidates_table.cell(0, 4).text = "% Valid Votes"

        for i, candidate in enumerate(candidates_slice):
            candidate_rank = candidate_rank + 1
            candidates_table.cell(i + 1, 0).text = str(candidate_rank)
            candidates_table.cell(i + 1, 1).text = candidate["candidate_name"]
            total_votes = candidate["total_votes"]
            candidates_table.cell(i + 1, 2).text = str(total_votes)
            valid_votes = candidate["valid_votes"]
            candidates_table.cell(i + 1, 3).text = str(valid_votes)
            if valid_votes == 0:
                candidates_table.cell(i + 1, 4).text = "0"
                continue
            candidates_table.cell(i + 1, 4).text = str(
                round(100 * total_votes / valid_votes, 2)
            )

    return


def create_results_power_point_headers(tally_id, filtered_electrol_races, qs):
    race_data_by_election_level_names = {
        f"{electrol_race.election_level}_{electrol_race.ballot_name}": {
            "election_level": electrol_race.election_level,
            "sub_race_type": electrol_race.ballot_name,
        }
        for electrol_race in filtered_electrol_races
        if electrol_race_has_results(qs, electrol_race)
    }
    tally_stations_qs = Station.objects.filter(tally_id=tally_id)
    stations_by_id = {station.id: station for station in tally_stations_qs}
    # Calculate voters in counted stations and turnout percentage
    for race_type_obj in race_data_by_election_level_names.values():
        # race_type_obj =\
        # race_data_by_election_level_names.get(election_level_name)
        sub_race_type = race_type_obj.get("sub_race_type")
        election_level_name = race_type_obj.get("election_level")
        voters = (
            Result.objects.filter(
                result_form__tally__id=tally_id,
                result_form__ballot__electrol_race__ballot_name=sub_race_type,
                result_form__form_state=FormState.ARCHIVED,
                entry_version=EntryVersion.FINAL,
                active=True,
            )
            .aggregate(race_voters=Coalesce(Sum("votes"), 0))
            .get("race_voters")
        )
        race_type_obj["voters_in_counted_stations"] = voters
        # Calculate voters in counted stations
        qs = tally_stations_qs.filter(
            center__resultform__ballot__electrol_race__election_level=election_level_name,
            center__resultform__ballot__electrol_race__ballot_name=sub_race_type,
        )
        race_type_obj["stations_expected"] = qs.distinct(
            "station_number", "center", "tally"
        ).count()

        station_ids_by_races = (
            qs.filter(
                center__resultform__form_state=FormState.ARCHIVED,
            )
            .annotate(
                race=F(
                    "center__resultform__ballot__electrol_race__election_level"
                )
            )
            .values("id")
            .annotate(
                races=ArrayAgg("race", distinct=True),
            )
        )
        stations_processed = 0
        registrants_in_processed_stations = 0
        for station in station_ids_by_races:
            station_obj = stations_by_id.get(station.get("id"))
            if station_obj is None:
                continue
            # Calculate stations processed and total registrants
            form_states = (
                ResultForm.objects.filter(
                    tally__id=tally_id,
                    center__resultform__ballot__electrol_race__election_level=election_level_name,
                    center__resultform__ballot__electrol_race__ballot_name=sub_race_type,
                    center__stations__id=station.get("id"),
                    station_number=station_obj.station_number,
                )
                .values_list("form_state", flat=True)
                .distinct()
            )

            station_is_processed = (
                form_states.count() == 1
                and form_states[0] == FormState.ARCHIVED
            )
            if station_is_processed is False:
                race_type_obj["stations_processed"] = 0
                race_type_obj["registrants_in_processed_stations"] = 0
                race_type_obj["percentage_of_stations_processed"] = 0
                race_type_obj["percentage_turnout"] = 0
                race_type_obj["results_status"] = "Partial"
                continue

            stations_processed += 1
            registrants_in_processed_stations += station_obj.registrants

        # Calculate turnout percentage
        if stations_processed != 0:
            race_type_obj["stations_processed"] = stations_processed
            race_type_obj["registrants_in_processed_stations"] = (
                registrants_in_processed_stations
            )
            race_type_obj["percentage_of_stations_processed"] = round(
                100 * stations_processed / race_type_obj["stations_expected"],
                2,
            )
            race_type_obj["results_status"] = (
                "Final"
                if race_type_obj["percentage_of_stations_processed"] >= 100.0
                else "Partial"
            )
            race_type_obj["percentage_turnout"] = round(
                100 * voters / registrants_in_processed_stations, 2
            )

    return race_data_by_election_level_names


def get_results(request):
    """
    Builds a json object of candidates results.

    :param request: The request object containing the tally id.

    returns: A JSON response of candidates results
    """
    tally_id = json.loads(request.GET.get("data")).get("tally_id")
    qs = Result.objects.select_related(
        "candidate",
        "result_form__center__sub_constituency",
        "result_form__ballot__electrol_race",
        "result_form__office",
    ).filter(
        result_form__tally__id=tally_id,
        result_form__form_state=FormState.ARCHIVED,
        entry_version=EntryVersion.FINAL,
        result_form__ballot__available_for_release=True,
        active=True,
    )
    results_metadata_queryset = qs
    qs = get_filtered_candidate_votes(
        tally_id,
        qs,
        data=None,
    )
    races_ids_to_total_votes_mapping = {}
    for electral_race_id in ElectrolRace.objects.filter(
        tally__id=tally_id
    ).values_list("id", flat=True):
        races_ids_to_total_votes_mapping[electral_race_id] = sum(
            [
                q.get("total_votes")
                for q in qs
                if q.get("electrol_race_id") == electral_race_id
            ]
        )

    candidate_data = []
    for candidate_result in qs:
        candidate_result.update(
            {
                "candidate_id": candidate_result.get("candidate_number"),
                "valid_votes": races_ids_to_total_votes_mapping.get(
                    candidate_result.get("electrol_race_id")
                ),
                "metadata": [
                    {
                        "barcode": rs.result_form.barcode,
                        "gender": rs.result_form.gender.name,
                        "station_number": rs.result_form.station_number,
                        "center_code": rs.result_form.center.code,
                        "center_name": rs.result_form.center.name,
                        "office_name": rs.result_form.office.name,
                        "office_number": rs.result_form.office.number,
                        "sub_con_name": \
                        rs.result_form.center.sub_constituency.name,
                        "sub_con_code": \
                        rs.result_form.center.sub_constituency.code,
                    }
                    for rs in results_metadata_queryset.filter(
                        candidate__candidate_id=candidate_result.get(
                            "candidate_number"
                        )
                    )
                ],
            }
        )
        candidate_data.append(candidate_result)

    sorted_results_report = sorted(
        candidate_data, key=lambda x: -x["total_votes"]
    )

    return JsonResponse(
        data={"data": sorted_results_report, "created_at": timezone.now()},
        safe=False,
    )


def get_centers_by_municipalities_results(request):
    """
    Builds a An Array of Dictionaries contain candidates total votes grouped
    by center code, sub race and sub constituency code.

    :param request: The request object containing the tally id.

    returns: A JSON response
    """
    tally_id = json.loads(request.GET.get("data")).get("tally_id")
    qs = Result.objects.select_related(
        "candidate",
        "result_form__center__sub_constituency",
        "result_form__ballot__electrol_race",
        "result_form__office",
    ).filter(
        result_form__tally__id=tally_id,
        result_form__form_state=FormState.ARCHIVED,
        entry_version=EntryVersion.FINAL,
        result_form__ballot__available_for_release=True,
        active=True,
    )
    data = (
        qs.annotate(
            code=F("result_form__center__code"),
            sub_race=F("result_form__ballot__electrol_race__ballot_name"),
            sub_con_code=F("result_form__center__sub_constituency__code"),
        )
        .values("code", "sub_race", "sub_con_code")
        .annotate(total_votes=Sum("votes"))
    )
    sorted_data = sorted(data, key=lambda x: -x["total_votes"])

    return JsonResponse(
        data={"data": sorted_data, "created_at": timezone.now()}, safe=False
    )


def get_centers_by_municipalities_candidates_results(request):
    """
    Builds a An Array of Dictionaries containing each candidates total votes
    grouped by center code, sub race and sub constituency code.

    :param request: The request object containing the tally id.

    returns: A JSON response
    """
    tally_id = json.loads(request.GET.get("data")).get("tally_id")
    qs = Result.objects.select_related(
        "candidate",
        "result_form__center__sub_constituency",
        "result_form__ballot__electrol_race",
        "result_form__office",
    ).filter(
        result_form__tally__id=tally_id,
        result_form__form_state=FormState.ARCHIVED,
        entry_version=EntryVersion.FINAL,
        result_form__ballot__available_for_release=True,
        active=True,
    )
    data = (
        qs.annotate(
            code=F("result_form__center__code"),
            sub_race=F("result_form__ballot__electrol_race__ballot_name"),
            sub_con_code=F("result_form__center__sub_constituency__code"),
            candidate_number=F("candidate__candidate_id"),
            candidate_name=F("candidate__full_name"),
        )
        .values(
            "code",
            "sub_race",
            "sub_con_code",
            "candidate_number",
            "candidate_name",
        )
        .annotate(total_votes=Sum("votes"))
    )
    sorted_data = sorted(data, key=lambda x: -x["total_votes"])

    return JsonResponse(
        data={"data": sorted_data, "created_at": timezone.now()}, safe=False
    )


def get_centers_stations_by_municipalities_candidates_results(request):
    """
    Builds a An Array of Dictionaries containing each candidates total votes
    grouped by center code, station number, sub race and sub constituency code.

    :param request: The request object containing the tally id.

    returns: A JSON response
    """
    tally_id = json.loads(request.GET.get("data")).get("tally_id")
    qs = Result.objects.select_related(
        "candidate",
        "result_form__center__sub_constituency",
        "result_form__ballot__electrol_race",
        "result_form__office",
    ).filter(
        result_form__tally__id=tally_id,
        result_form__form_state=FormState.ARCHIVED,
        entry_version=EntryVersion.FINAL,
        result_form__ballot__available_for_release=True,
        active=True,
    )
    data = (
        qs.annotate(
            code=F("result_form__center__code"),
            station_number=F("result_form__station_number"),
            sub_race=F("result_form__ballot__electrol_race__ballot_name"),
            sub_con_code=F("result_form__center__sub_constituency__code"),
            candidate_number=F("candidate__candidate_id"),
            candidate_name=F("candidate__full_name"),
        )
        .values(
            "code",
            "station_number",
            "sub_race",
            "sub_con_code",
            "candidate_number",
            "candidate_name",
        )
        .annotate(total_votes=Sum("votes"))
    )
    sorted_data = sorted(data, key=lambda x: -x["total_votes"])

    return JsonResponse(
        data={"data": sorted_data, "created_at": timezone.now()}, safe=False
    )


def get_sub_cons_list(request):
    """
    Builds a json object of sub constituencies.

    :param request: The request object containing the tally id.

    returns: A JSON response of sub constituencies results
    """
    tally_id = json.loads(request.GET.get("data")).get("tally_id")
    qs = (
        SubConstituency.objects.filter(tally__id=tally_id)
        .prefetch_related("ballots__electrol_race")
        .values(
            "code",
            "name",
            "field_office",
            "ballots__electrol_race__election_level",
            "ballots__electrol_race__ballot_name",
            "ballots__number",
        )
        .annotate(
            election_level=F("ballots__electrol_race__election_level"),
            sub_race=F("ballots__electrol_race__ballot_name"),
            ballot_number=F("ballots__number"),
        )
    )
    subs = [
        {
            "code": sub.get("code"),
            "name": sub.get("name"),
            "election_level": sub.get("election_level"),
            "sub_race": sub.get("sub_race"),
            "ballot_number": sub.get("ballot_number"),
        }
        for sub in qs
    ]

    return JsonResponse(
        data={"data": subs, "created_at": timezone.now()}, safe=False
    )


class SummaryReportDataView(
    LoginRequiredMixin,
    GroupRequiredMixin,
    TallyAccessMixin,
    BaseDatatableView,
):
    group_required = groups.TALLY_MANAGER
    model = ReconciliationForm
    columns = (
        "name",
        "number_valid_votes",
        "number_invalid_votes",
        "constituencies_ids",
        "sub_constituencies_ids",
        "actions",
    )

    def filter_queryset(self, qs):
        tally_id = self.kwargs.get("tally_id")
        region_id = self.kwargs.get("region_id")
        constituency_id = self.kwargs.get("constituency_id")
        data = self.request.POST.get("data")
        keyword = self.request.POST.get("search[value]")

        if data:
            qs = custom_queryset_filter(
                tally_id,
                qs,
                ast.literal_eval(data),
                report_type="summary",
                region_id=region_id,
                constituency_id=constituency_id,
            )
        else:
            qs = custom_queryset_filter(
                tally_id,
                qs,
                report_type="summary",
                region_id=region_id,
                constituency_id=constituency_id,
            )

        if keyword:
            qs = qs.filter(
                Q(name__icontains=keyword)
                | Q(total_number_of_registrants__contains=keyword)
                | Q(number_of_voters_voted__contains=keyword)
                | Q(male_voters__contains=keyword)
                | Q(female_voters__contains=keyword)
                | Q(turnout_percentage__contains=keyword)
            )
        return qs

    def render_column(self, row, column):
        tally_id = self.kwargs.get("tally_id")
        region_id = self.kwargs.get("region_id")
        constituency_id = self.kwargs.get("constituency_id")
        data = self.request.POST.get("data")
        administrative_area_child_report_name = _("Region Constituencies")
        url = reverse(
            "constituency-summary-report",
            kwargs={"tally_id": tally_id, "region_id": row["admin_area_id"]},
        )

        if region_id:
            administrative_area_child_report_name = _("Sub Constituencies")
            url = reverse(
                "sub-constituency-summary-report",
                kwargs={
                    "tally_id": tally_id,
                    "region_id": row["admin_area_id"],
                    "constituency_id": row["constituency_id"],
                },
            )

        if column == "name":
            return str(f'<td class="center">{row["name"]}</td>')
        elif column == "number_valid_votes":
            return str(f'<td class="center">{row["number_valid_votes"]}</td>')
        elif column == "number_invalid_votes":
            return str(
                f'<td class="center">{row["number_invalid_votes"]}</td>'
            )
        elif column == "constituencies_ids":
            disabled = "disabled" if region_id else ""
            region_cons_ids = []
            qs = Constituency.objects.filter(
                tally__id=tally_id, id__in=row["constituencies_ids"]
            ).values_list("id", "name", named=True)
            if data:
                region_cons_data = [
                    item
                    for item in ast.literal_eval(data)
                    if ast.literal_eval(item["region_id"])
                    == row["admin_area_id"]
                ]
                region_cons_ids = region_cons_data[0]["select_1_ids"]
            constituencies = build_select_options(qs, ids=region_cons_ids)
            return str(
                '<td class="center">'
                '<select style="min-width: 6em;"'
                f"{disabled}"
                ' id="select-1" multiple'
                " data-id="
                f"{row['admin_area_id']}"
                ">"
                f"{constituencies}"
                "</select>"
                "</td>"
            )
        elif column == "sub_constituencies_ids":
            disabled = "disabled" if constituency_id else ""
            region_sub_cons_ids = []
            qs = (
                SubConstituency.objects.annotate(sc_code=F("code"))
                .filter(
                    tally__id=tally_id, id__in=row["sub_constituencies_ids"]
                )
                .values_list("id", "sc_code", named=True)
            )
            qs = [(item.id, item.sc_code) for item in qs]
            if data:
                region_sub_cons_data = [
                    item
                    for item in ast.literal_eval(data)
                    if ast.literal_eval(item["region_id"])
                    == row["admin_area_id"]
                ]
                region_sub_cons_ids = region_sub_cons_data[0]["select_2_ids"]

            sub_constituencies = build_select_options(
                qs, ids=region_sub_cons_ids
            )
            return str(
                '<td class="center">'
                '<select style="min-width: 6em !important;"'
                f"{disabled}"
                ' id="select-2" multiple'
                " data-id="
                f"{row['admin_area_id']}"
                ">"
                f"{sub_constituencies}"
                "</select>"
                "</td>"
            )
        elif column == "actions":
            if constituency_id:
                return str(
                    '<button id="filter-report" disabled '
                    'class="btn btn-default btn-small">Submit</button>'
                )
            return str(
                "<a href="
                f"{url}"
                ' class="btn btn-default btn-small vertical-margin"> '
                f"{administrative_area_child_report_name}"
                "</a>"
                '<button id="filter-report" '
                'class="btn btn-default btn-small">Submit</button>'
            )
        else:
            return super(SummaryReportDataView, self).render_column(
                row, column
            )


class SummaryReportView(
    LoginRequiredMixin, GroupRequiredMixin, DataTablesMixin, TemplateView
):
    group_required = groups.TALLY_MANAGER
    model = ReconciliationForm
    template_name = "reports/summary_report.html"

    def get(self, request, *args, **kwargs):
        tally_id = kwargs.get("tally_id")
        region_id = kwargs.get("region_id")
        constituency_id = kwargs.get("constituency_id")

        try:
            region_name = (
                region_id
                and Region.objects.get(id=region_id, tally__id=tally_id).name
            )
        except Region.DoesNotExist:
            region_name = None

        try:
            constituency_name = (
                constituency_id
                and Constituency.objects.get(
                    id=constituency_id, tally__id=tally_id
                ).name
            )
        except Constituency.DoesNotExist:
            constituency_name = None

        return self.render_to_response(
            self.get_context_data(
                remote_url=reverse("summary-list-data", kwargs=kwargs),
                tally_id=tally_id,
                region_name=region_name,
                constituency_name=constituency_name,
            )
        )


class ProgressiveReportDataView(
    LoginRequiredMixin,
    GroupRequiredMixin,
    TallyAccessMixin,
    BaseDatatableView,
):
    group_required = groups.TALLY_MANAGER
    model = Result
    columns = (
        "admin_area_name",
        "total_candidates",
        "total_votes",
        "constituencies_ids",
        "sub_constituencies_ids",
        "actions",
    )

    def filter_queryset(self, qs):
        tally_id = self.kwargs.get("tally_id")
        region_id = self.kwargs.get("region_id")
        constituency_id = self.kwargs.get("constituency_id")
        data = self.request.POST.get("data")
        keyword = self.request.POST.get("search[value]")
        qs = qs.filter(
            result_form__tally__id=tally_id,
            result_form__form_state=FormState.ARCHIVED,
            entry_version=EntryVersion.FINAL,
            active=True,
        )

        if data:
            qs = generate_progressive_report_queryset(
                qs,
                ast.literal_eval(data),
                region_id=region_id,
                constituency_id=constituency_id,
            )
        else:
            qs = generate_progressive_report_queryset(
                qs, region_id=region_id, constituency_id=constituency_id
            )

        if keyword:
            qs = qs.filter(
                Q(admin_area_name__icontains=keyword)
                | Q(total_candidates=keyword)
                | Q(total_votes=keyword)
            )
        return qs

    def render_column(self, row, column):
        tally_id = self.kwargs.get("tally_id")
        region_id = self.kwargs.get("region_id")
        constituency_id = self.kwargs.get("constituency_id")
        data = self.request.POST.get("data")
        child_report_button_text = None
        child_report_url = None
        votes_per_candidate_url = None
        votes_per_candidate_button_text = None

        if region_id and not constituency_id:
            reverse_url = "sub-cons-progressive-report-list"
            child_report_button_text = _(
                "Sub Constituencies votes per candidate"
            )
            child_report_url = reverse(
                reverse_url,
                kwargs={
                    "tally_id": tally_id,
                    "region_id": row["region_id"],
                    "constituency_id": row["constituency_id"],
                },
            )

            reverse_url = "constituency-votes-per-candidate"
            votes_per_candidate_button_text = _(
                "Constituency votes per candidate"
            )
            votes_per_candidate_url = reverse(
                reverse_url,
                kwargs={
                    "tally_id": tally_id,
                    "region_id": row["region_id"],
                    "constituency_id": row["constituency_id"],
                    "report_type": "votes-per-candidate-report",
                },
            )
        elif region_id and constituency_id:
            reverse_url = "sub-constituency-votes-per-candidate"
            child_report_button_text = _(
                "Sub Constituency votes per candidate"
            )
            child_report_url = reverse(
                reverse_url,
                kwargs={
                    "tally_id": tally_id,
                    "region_id": row["region_id"],
                    "constituency_id": row["constituency_id"],
                    "sub_constituency_id": row["sub_constituency_id"],
                    "report_type": "votes-per-candidate-report",
                },
            )

            reverse_url = "sub-constituency-votes-per-candidate"
            votes_per_candidate_button_text = _(
                "Sub Constituency candidates list by ballot order"
            )
            votes_per_candidate_url = reverse(
                reverse_url,
                kwargs={
                    "tally_id": tally_id,
                    "region_id": row["region_id"],
                    "constituency_id": row["constituency_id"],
                    "sub_constituency_id": row["sub_constituency_id"],
                    "report_type": "candidate-list-sorted-by-ballots-number",
                },
            )
        else:
            reverse_url = "cons-progressive-report-list"
            child_report_button_text = _(
                "Region Constituencies Progressive Report"
            )
            child_report_url = reverse(
                reverse_url,
                kwargs={"tally_id": tally_id, "region_id": row["region_id"]},
            )

            reverse_url = "region-votes-per-candidate"
            votes_per_candidate_button_text = _("Region votes per candidate")
            votes_per_candidate_url = reverse(
                reverse_url,
                kwargs={
                    "tally_id": tally_id,
                    "region_id": row["region_id"],
                    "report_type": "votes-per-candidate-report",
                },
            )

        if column == "admin_area_name":
            return str(f'<td class="center">{row["admin_area_name"]}</td>')
        elif column == "total_candidates":
            return str(f'<td class="center">{row["total_candidates"]}</td>')
        elif column == "total_votes":
            return str(f'<td class="center">{row["total_votes"]}</td>')
        elif column == "constituencies_ids":
            disabled = "disabled" if region_id else ""
            region_cons_ids = []
            qs = Constituency.objects.filter(
                tally__id=tally_id, id__in=row["constituencies_ids"]
            ).values_list("id", "name", named=True)
            if data:
                region_cons_data = [
                    item
                    for item in ast.literal_eval(data)
                    if ast.literal_eval(item["region_id"])
                    == row["admin_area_id"]
                ]
                region_cons_ids = region_cons_data[0]["select_1_ids"]
            constituencies = build_select_options(qs, ids=region_cons_ids)
            return str(
                '<td class="center">'
                '<select style="min-width: 6em;"'
                f"{disabled}"
                ' id="select-1" multiple'
                " data-id="
                f"{row['admin_area_id']}"
                ">"
                f"{constituencies}"
                "</select>"
                "</td>"
            )
        elif column == "sub_constituencies_ids":
            disabled = "disabled" if constituency_id else ""
            region_sub_cons_ids = []
            qs = (
                SubConstituency.objects.annotate(sc_code=F("code"))
                .filter(
                    tally__id=tally_id, id__in=row["sub_constituencies_ids"]
                )
                .values_list("id", "sc_code", named=True)
            )
            qs = [(item.id, item.sc_code) for item in qs]
            if data:
                region_sub_cons_data = [
                    item
                    for item in ast.literal_eval(data)
                    if ast.literal_eval(item["region_id"])
                    == row["admin_area_id"]
                ]
                region_sub_cons_ids = region_sub_cons_data[0]["select_2_ids"]

            sub_constituencies = build_select_options(
                qs, ids=region_sub_cons_ids
            )
            return str(
                '<td class="center">'
                '<select style="min-width: 6em;"'
                f"{disabled}"
                ' id="select-2" multiple'
                " data-id="
                f"{row['admin_area_id']}"
                ">"
                f"{sub_constituencies}"
                "</select>"
                "</td>"
            )
        elif column == "actions":
            child_report_link = str(
                "<a href="
                f"{child_report_url}"
                ' class="btn btn-default btn-small vertical-margin"> '
                f"{child_report_button_text}"
                "</a>"
            )
            votes_per_candidate_link = str(
                "<a href="
                f"{votes_per_candidate_url}"
                ' class="btn btn-default btn-small vertical-margin"> '
                f"{votes_per_candidate_button_text}"
                "</a>"
            )
            filter_button = str(
                '<button id="filter-report" '
                'class="btn btn-default btn-small">Submit</button>'
            )
            return child_report_link + votes_per_candidate_link + filter_button
        else:
            return super(ProgressiveReportDataView, self).render_column(
                row, column
            )


class ProgressiveReportView(
    LoginRequiredMixin, GroupRequiredMixin, DataTablesMixin, TemplateView
):
    group_required = groups.TALLY_MANAGER
    model = Result
    template_name = "reports/progressive_report.html"

    def get(self, request, *args, **kwargs):
        tally_id = kwargs.get("tally_id")
        region_id = kwargs.get("region_id")
        constituency_id = kwargs.get("constituency_id")

        try:
            region_name = (
                region_id
                and Region.objects.get(id=region_id, tally__id=tally_id).name
            )
        except Region.DoesNotExist:
            region_name = None

        try:
            constituency_name = (
                constituency_id
                and Constituency.objects.get(
                    id=constituency_id, tally__id=tally_id
                ).name
            )
        except Constituency.DoesNotExist:
            constituency_name = None

        return self.render_to_response(
            self.get_context_data(
                remote_url=reverse(
                    "progressive-report-list-data", kwargs=kwargs
                ),
                tally_id=tally_id,
                region_name=region_name,
                constituency_name=constituency_name,
            )
        )


class DiscrepancyReportDataView(
    LoginRequiredMixin,
    GroupRequiredMixin,
    TallyAccessMixin,
    BaseDatatableView,
):
    group_required = groups.TALLY_MANAGER
    model = Station
    columns = (
        "admin_area_name",
        "number_of_centers",
        "number_of_stations",
        "station_ids",
        "center_ids",
        "actions",
    )

    def filter_queryset(self, qs):
        tally_id = self.kwargs.get("tally_id")
        region_id = self.kwargs.get("region_id")
        constituency_id = self.kwargs.get("constituency_id")
        data = self.request.POST.get("data")
        keyword = self.request.POST.get("search[value]")
        report_name = self.kwargs.get("report_name")
        stations_centers_under_process_audit = report_types[3]

        if report_name == stations_centers_under_process_audit:
            qs = ResultForm.objects.filter(
                tally__id=tally_id, form_state=FormState.AUDIT
            )

        if data:
            qs = stations_and_centers_queryset(
                tally_id,
                qs,
                ast.literal_eval(data),
                report_type=report_name,
                region_id=region_id,
                constituency_id=constituency_id,
            )
        else:
            qs = stations_and_centers_queryset(
                tally_id,
                qs,
                report_type=report_name,
                region_id=region_id,
                constituency_id=constituency_id,
            )

        if keyword:
            qs = qs.filter(
                Q(admin_area_name__icontains=keyword)
                | Q(number_of_centers__contains=keyword)
                | Q(number_of_stations__contains=keyword)
            )
        return qs

    def render_column(self, row, column):
        tally_id = self.kwargs.get("tally_id")
        region_id = self.kwargs.get("region_id")
        constituency_id = self.kwargs.get("constituency_id")
        data = self.request.POST.get("data")
        child_report_button_text = None
        child_report_url = None
        station_and_centers_list_url = None
        station_and_centers_list_button_text = None

        report_name = self.kwargs.get("report_name")
        stations_centers_under_process_audit = report_types[3]
        stations_centers_under_investigation = report_types[4]
        stations_centers_excluded_after_investigation = report_types[5]

        if report_name == stations_centers_under_investigation:
            if region_id and not constituency_id:
                reverse_url = (
                    "sub-cons-stations-and-centers-under-investigation"
                )
                button_text = (
                    "Sub Constituency Station and Centers under investigation"
                )
                child_report_button_text = _(button_text)
                child_report_url = reverse(
                    reverse_url,
                    kwargs={
                        "tally_id": tally_id,
                        "region_id": row["region_id"],
                        "constituency_id": row["constituency_id"],
                        "report_name": \
                        "stations-and-centers-under-investigation-list",
                    },
                )

                reverse_url = "constituency-discrepancy-report"
                station_and_centers_list_button_text = _(
                    "Constituency Centers and Stations under investigation"
                )
                station_and_centers_list_url = reverse(
                    reverse_url,
                    kwargs={
                        "tally_id": tally_id,
                        "region_id": row["region_id"],
                        "constituency_id": row["constituency_id"],
                        "report_type": \
                        "centers-and-stations-under-investigation",
                    },
                )
            elif region_id and constituency_id:
                reverse_url = "sub-constituency-discrepancy-report"
                button_text_1 = "Sub Constituency Stations and Centers"
                button_text_2 = " under investigation"
                station_and_centers_list_button_text = _(
                    "{}{}".format(button_text_1, button_text_2)
                )
                station_and_centers_list_url = reverse(
                    reverse_url,
                    kwargs={
                        "tally_id": tally_id,
                        "region_id": row["region_id"],
                        "constituency_id": row["constituency_id"],
                        "sub_constituency_id": row["sub_constituency_id"],
                        "report_type": \
                        "centers-and-stations-under-investigation",
                    },
                )
            else:
                reverse_url = "cons-stations-and-centers-under-investigation"
                child_report_button_text = _(
                    "Region Constituencies under Investigation"
                )
                report_name = "stations-and-centers-under-investigation-list"
                child_report_url = reverse(
                    reverse_url,
                    kwargs={
                        "tally_id": tally_id,
                        "region_id": row["region_id"],
                        "report_name": report_name,
                    },
                )

                reverse_url = "regions-discrepancy-report"
                station_and_centers_list_button_text = _(
                    "Region Centers and Stations under Investigation"
                )
                report_type = "centers-and-stations-under-investigation"
                station_and_centers_list_url = reverse(
                    reverse_url,
                    kwargs={
                        "tally_id": tally_id,
                        "region_id": row["region_id"],
                        "report_type": report_type,
                    },
                )

        elif report_name == stations_centers_excluded_after_investigation:
            if region_id and not constituency_id:
                url_part_1 = "sub-cons-stations-and-centers-excluded"
                url_part_2 = "-after-investigation"
                reverse_url = f"{url_part_1}{url_part_2}"
                button_text_1 = "Sub Constituency Station and Centers"
                button_text_2 = " excluded after investigation"
                child_report_button_text = _(
                    "{}{}".format(button_text_1, button_text_2)
                )
                report_name = (
                    "stations-and-centers-excluded-after-investigation-list"
                )
                child_report_url = reverse(
                    reverse_url,
                    kwargs={
                        "tally_id": tally_id,
                        "region_id": row["region_id"],
                        "constituency_id": row["constituency_id"],
                        "report_name": report_name,
                    },
                )

                reverse_url = "constituency-discrepancy-report"
                button_text_1 = "Constituency Centers and Stations excluded"
                button_text_2 = " after investigation"
                station_and_centers_list_button_text = _(
                    "{}{}".format(button_text_1, button_text_2)
                )
                report_type = (
                    "centers-and-stations-excluded-after-investigation"
                )
                station_and_centers_list_url = reverse(
                    reverse_url,
                    kwargs={
                        "tally_id": tally_id,
                        "region_id": row["region_id"],
                        "constituency_id": row["constituency_id"],
                        "report_type": report_type,
                    },
                )
            elif region_id and constituency_id:
                reverse_url = "sub-constituency-discrepancy-report"
                button_text_1 = (
                    "Sub Constituency Centers and Stations excluded "
                )
                button_text_2 = "after investigation"
                station_and_centers_list_button_text = _(
                    "{}{}".format(button_text_1, button_text_2)
                )
                report_type = (
                    "centers-and-stations-excluded-after-investigation"
                )
                station_and_centers_list_url = reverse(
                    reverse_url,
                    kwargs={
                        "tally_id": tally_id,
                        "region_id": row["region_id"],
                        "constituency_id": row["constituency_id"],
                        "sub_constituency_id": row["sub_constituency_id"],
                        "report_type": report_type,
                    },
                )
            else:
                reverse_url = (
                    "cons-stations-and-centers-excluded-after-investigation"
                )
                child_report_button_text = _(
                 "Region Constituencies excluded after investigation"
                )
                report_name = (
                    "stations-and-centers-excluded-after-investigation-list"
                )
                child_report_url = reverse(
                    reverse_url,
                    kwargs={
                        "tally_id": tally_id,
                        "region_id": row["region_id"],
                        "report_name": report_name,
                    },
                )

                reverse_url = "regions-discrepancy-report"
                button_text = (
                    "Region Centers and Stations excluded after investigation"
                )
                station_and_centers_list_button_text = _(button_text)
                report_type = (
                    "centers-and-stations-excluded-after-investigation"
                )
                station_and_centers_list_url = reverse(
                    reverse_url,
                    kwargs={
                        "tally_id": tally_id,
                        "region_id": row["region_id"],
                        "report_type": report_type,
                    },
                )
        elif report_name == stations_centers_under_process_audit:
            if region_id and not constituency_id:
                reverse_url = (
                    "sub-cons-stations-and-centers-under-process-audit-list"
                )
                child_report_button_text = _(
                    "Sub Constituencies with Station and Centers in Audit"
                )
                child_report_url = reverse(
                    reverse_url,
                    kwargs={
                        "tally_id": tally_id,
                        "region_id": row["region_id"],
                        "constituency_id": row["constituency_id"],
                        "report_name": \
                        "stations-and-centers-under-process-audit-list",
                    },
                )

                reverse_url = "constituency-discrepancy-report"
                station_and_centers_list_button_text = _(
                    "Constituency Centers and Stations in Audit"
                )
                station_and_centers_list_url = reverse(
                    reverse_url,
                    kwargs={
                        "tally_id": tally_id,
                        "region_id": row["region_id"],
                        "constituency_id": row["constituency_id"],
                        "report_type": "centers-and-stations-in-audit-report",
                    },
                )
            elif region_id and constituency_id:
                reverse_url = "sub-constituency-discrepancy-report"
                station_and_centers_list_button_text = _(
                    "Sub Constituency Stations and Centers in Audit"
                )
                station_and_centers_list_url = reverse(
                    reverse_url,
                    kwargs={
                        "tally_id": tally_id,
                        "region_id": row["region_id"],
                        "constituency_id": row["constituency_id"],
                        "sub_constituency_id": row["sub_constituency_id"],
                        "report_type": "centers-and-stations-in-audit-report",
                    },
                )
            else:
                reverse_url = (
                    "cons-stations-and-centers-under-process-audit-list"
                )
                button_text = (
                    "Region Constituencies with Stations and Centers in Audit"
                )
                child_report_button_text = _(button_text)
                child_report_url = reverse(
                    reverse_url,
                    kwargs={
                        "tally_id": tally_id,
                        "region_id": row["region_id"],
                        "report_name": \
                        "stations-and-centers-under-process-audit-list",
                    },
                )

                reverse_url = "regions-discrepancy-report"
                station_and_centers_list_button_text = _(
                    "Region Centers and Stations under process Audit"
                )
                station_and_centers_list_url = reverse(
                    reverse_url,
                    kwargs={
                        "tally_id": tally_id,
                        "region_id": row["region_id"],
                        "report_type": "centers-and-stations-in-audit-report",
                    },
                )

        if column == "admin_area_name":
            return str(f'<td class="center">{row["admin_area_name"]}</td>')
        elif column == "number_of_centers":
            return str(f'<td class="center">{row["number_of_centers"]}</td>')
        elif column == "number_of_stations":
            return str(f'<td class="center">{row["number_of_stations"]}</td>')
        elif column == "station_ids":
            region_station_ids = []
            station_ids = row["station_ids"]
            if report_name == stations_centers_under_process_audit:
                station_ids = list(
                    Station.objects.filter(
                        station_number__in=row["station_ids"],
                        center__id__in=row["center_ids"],
                        tally__id=tally_id,
                    )
                    .distinct("station_number")
                    .values_list("id", flat=True)
                )
            qs = (
                Station.objects.annotate(name=F("station_number"))
                .filter(tally__id=tally_id, id__in=station_ids)
                .values_list("id", "name", named=True)
            )
            if data:
                region_stations_data = [
                    item
                    for item in ast.literal_eval(data)
                    if ast.literal_eval(item["region_id"])
                    == row["admin_area_id"]
                ]
                region_station_ids = region_stations_data[0]["select_1_ids"]
            stations = build_select_options(qs, ids=region_station_ids)
            disabled = "disabled" if not len(stations) else ""

            return str(
                '<td class="center">'
                '<select style="min-width: 6em;" '
                f"{disabled}"
                ' id="select-1" multiple'
                " data-id="
                f"{row['admin_area_id']}"
                ">"
                f"{stations}"
                "</select>"
                "</td>"
            )
        elif column == "center_ids":
            region_center_ids = []
            qs = Center.objects.filter(
                tally__id=tally_id, id__in=row["center_ids"]
            ).values_list("id", "name", named=True)

            if data:
                region_centers_data = [
                    item
                    for item in ast.literal_eval(data)
                    if ast.literal_eval(item["region_id"])
                    == row["admin_area_id"]
                ]
                region_center_ids = region_centers_data[0]["select_2_ids"]

            centers = build_select_options(qs, ids=region_center_ids)
            disabled = "disabled" if not len(centers) else ""

            return str(
                '<td class="center">'
                '<select style="min-width: 6em;" '
                f"{disabled}"
                ' id="select-2" multiple'
                " data-id="
                f"{row['admin_area_id']}"
                ">"
                f"{centers}"
                "</select>"
                "</td>"
            )
        elif column == "actions":
            child_report_link = ""
            station_and_centers_list_link = ""
            if child_report_url:
                child_report_link = str(
                    "<a href="
                    f"{child_report_url}"
                    ' class="btn btn-default btn-small vertical-margin"> '
                    f"{child_report_button_text}"
                    "</a>"
                )
            if station_and_centers_list_url:
                station_and_centers_list_link = str(
                    "<a href="
                    f"{station_and_centers_list_url}"
                    ' class="btn btn-default btn-small vertical-margin"> '
                    f"{station_and_centers_list_button_text}"
                    "</a>"
                )
            filter_button = str(
                '<button id="filter-report" '
                'class="btn btn-default btn-small">Submit</button>'
            )
            return (
                child_report_link
                + station_and_centers_list_link
                + filter_button
            )
        else:
            return super(DiscrepancyReportDataView, self).render_column(
                row, column
            )


class DiscrepancyReportView(
    LoginRequiredMixin, GroupRequiredMixin, DataTablesMixin, TemplateView
):
    group_required = groups.TALLY_MANAGER
    model = Station
    template_name = "reports/process_discrepancy_report.html"

    def get(self, request, *args, **kwargs):
        tally_id = kwargs.get("tally_id")
        region_id = kwargs.get("region_id")
        constituency_id = kwargs.get("constituency_id")
        report_type = None
        url = None

        report_name = kwargs.get("report_name")
        stations_centers_under_process_audit = report_types[3]
        stations_centers_under_investigation = report_types[4]
        stations_centers_excluded_after_investigation = report_types[5]
        if report_name == stations_centers_under_process_audit:
            report_type = _("Stations and Centers under process audit")
            url = "stations-and-centers-under-process-audit-list-data"
        elif report_name == stations_centers_under_investigation:
            report_type = _("Stations and Centers under investigation")
            url = "stations-and-centers-under-investigation-list-data"
        elif report_name == stations_centers_excluded_after_investigation:
            report_type = _(
                "Stations and Centers excluded after investigation"
            )
            url = "stations-and-centers-excluded-after-investigation-data"

        try:
            region_name = (
                region_id
                and Region.objects.get(id=region_id, tally__id=tally_id).name
            )
        except Region.DoesNotExist:
            region_name = None

        try:
            constituency_name = (
                constituency_id
                and Constituency.objects.get(
                    id=constituency_id, tally__id=tally_id
                ).name
            )
        except Constituency.DoesNotExist:
            constituency_name = None

        return self.render_to_response(
            self.get_context_data(
                remote_url=reverse(url, kwargs=kwargs),
                tally_id=tally_id,
                region_name=region_name,
                constituency_name=constituency_name,
                report_type=report_type,
            )
        )


def generate_report(
    tally_id,
    report_column_name,
    report_type_name,
    region_id=None,
    constituency_id=None,
):
    """
    Genarate report by using the final reconciliation form to get voter stats.

    :param tally_id: The reconciliation forms tally.
    :param report_column_name: The result form report column name.
    :param region_id: The region id for filtering the recon forms.
    :param constituency_id: The constituency id for filtering the recon forms.
    :param report_type_name: The report type name to generate.

    returns: The turnout report grouped by the report column name.
    """
    turnout_report_type_name = report_types[1]
    summary_report_type_name = report_types[2]

    qs = ReconciliationForm.objects.get_registrants_and_votes_type().filter(
        result_form__tally__id=tally_id,
        result_form__form_state=FormState.ARCHIVED,
        entry_version=EntryVersion.FINAL,
    )
    if region_id:
        qs = qs.filter(result_form__office__region__id=region_id)

    if constituency_id:
        qs = qs.filter(result_form__center__constituency__id=constituency_id)
    qs = qs.annotate(
        name=F(report_column_name),
        admin_area_id=F("result_form__office__region__id"),
    ).values(
        "name",
        "admin_area_id",
    )

    if region_id:
        qs = qs.annotate(
            constituency_id=F("result_form__center__constituency__id"),
        )

    if report_type_name == turnout_report_type_name:
        qs = (
            qs.annotate(number_of_voters_voted=Sum("number_valid_votes"))
            .annotate(
                total_number_of_registrants=Sum(
                    "result_form__center__stations__registrants"
                )
            )
            .annotate(
                total_number_of_ballots_used=Sum(
                    ExpressionWrapper(
                        F("number_valid_votes")
                        + F("number_invalid_votes"),
                        output_field=IntegerField(),
                    )
                )
            )
            .annotate(
                turnout_percentage=ExpressionWrapper(
                    V(100)
                    * F("total_number_of_ballots_used")
                    / F("total_number_of_registrants"),
                    output_field=IntegerField(),
                )
            )
            .annotate(
                male_voters=Coalesce(
                    Sum("number_valid_votes", filter=Q(voters_gender_type=0)),
                    V(0),
                )
            )
            .annotate(
                female_voters=Coalesce(
                    Sum("number_valid_votes", filter=Q(voters_gender_type=1)),
                    V(0),
                )
            )
        )

    if report_type_name == summary_report_type_name:
        qs = (
            qs.annotate(number_valid_votes=Sum("number_valid_votes"))
            .annotate(number_invalid_votes=Sum("number_invalid_votes"))
        )

    return qs


class RegionsReportsView(
    LoginRequiredMixin, GroupRequiredMixin, TemplateView
):
    group_required = groups.TALLY_MANAGER
    template_name = "reports/administrative_areas_reports.html"

    def get(self, request, *args, **kwargs):
        tally_id = kwargs["tally_id"]
        report_type_ = kwargs.get("report_type")
        region_id = kwargs.get("region_id")
        column_name = "result_form__office__region__name"
        qs = Station.objects.filter(tally__id=tally_id)

        progressive_report = generate_progressive_report(
            tally_id=tally_id,
            report_column_name=column_name,
        )

        centers_stations_in_audit = stations_and_centers_queryset(
            tally_id=tally_id,
            qs=ResultForm.objects.filter(
                tally__id=tally_id, form_state=FormState.AUDIT
            ),
            report_type=report_types[3],
        )

        centers_stations_under_invg = stations_and_centers_queryset(
            tally_id=tally_id, qs=qs, report_type=report_types[4]
        )

        centers_stations_ex_after_invg = stations_and_centers_queryset(
            tally_id=tally_id, qs=qs, report_type=report_types[5]
        )

        station_id_query = Subquery(
            Station.objects.filter(
                tally__id=tally_id,
                center__code=OuterRef("center__code"),
                station_number=OuterRef("station_number"),
            ).values("id")[:1],
            output_field=IntegerField(),
        )

        if report_type_ in [
            "centers-and-stations-in-audit-report",
            "centers-and-stations-under-investigation",
            "centers-and-stations-excluded-after-investigation",
        ]:
            if report_type_ == "centers-and-stations-in-audit-report":
                self.request.session["station_ids"] = list(
                    centers_stations_in_audit.filter(
                        center__office__region__id=region_id
                    )
                    .annotate(station_id=station_id_query)
                    .values_list("station_id", flat=True)
                )

            if report_type_ == "centers-and-stations-under-investigation":
                self.request.session["station_ids"] = list(
                    centers_stations_under_invg.filter(
                        center__office__region__id=region_id
                    )
                    .annotate(station_id=station_id_query)
                    .values_list("station_id", flat=True)
                )

            if (
                report_type_
                == "centers-and-stations-excluded-after-investigation"
            ):
                self.request.session["station_ids"] = list(
                    centers_stations_ex_after_invg.filter(
                        center__office__region__id=region_id
                    )
                    .annotate(station_id=station_id_query)
                    .values_list("station_id", flat=True)
                )

            return redirect(
                "center-and-stations-list",
                tally_id=tally_id,
                region_id=region_id,
            )

        if report_type_ == "votes-per-candidate-report":
            self.request.session["result_ids"] = list(
                progressive_report.filter(
                    result_form__center__office__region__id=region_id
                ).values_list("id", flat=True)
            )
            self.request.session["ballot_report"] = False

            return redirect(
                "candidate-list-by-votes",
                tally_id=tally_id,
                region_id=region_id,
            )

        return self.render_to_response(
            self.get_context_data(
                tally_id=tally_id,
                administrative_area_name=None,
                region_name=None,
                constituency_name=None,
                report_name=_("Region"),
                administrative_area_child_report_name=_(
                    "Region Constituencies"
                ),
                turn_out_report_download_url="regions-turnout-csv",
                summary_report_download_url="regions-summary-csv",
                centers_stations_under_invg=centers_stations_under_invg,
                centers_stations_ex_after_invg=centers_stations_ex_after_invg,
                regions_report_url="regions-discrepancy-report",
                child_summary_report_url="constituency-summary-report",
                child_discrepancy_report_url="constituency-discrepancy-report",
                child_progressive_report_url="constituency-progressive-report",
                admin_area_votes_per_candidate_report_url="region-votes-per-candidate",
            )
        )


class ConstituencyReportsView(
    LoginRequiredMixin, GroupRequiredMixin, TemplateView
):
    group_required = groups.TALLY_MANAGER

    def get(self, request, *args, **kwargs):
        tally_id = kwargs["tally_id"]
        region_id = kwargs["region_id"]
        report_type = kwargs.get("report_type")
        constituency_id = kwargs.get("constituency_id")

        region_name = (
            Region.objects.get(id=region_id, tally__id=tally_id).name
            if region_id
            else None
        )
        column_name = "result_form__center__constituency__name"
        progressive_report = generate_progressive_report(
            tally_id=tally_id,
            report_column_name=column_name,
            region_id=region_id,
        )

        qs = Station.objects.filter(tally__id=tally_id)

        centers_stations_in_audit = stations_and_centers_queryset(
            tally_id=tally_id,
            qs=ResultForm.objects.filter(
                tally__id=tally_id, form_state=FormState.AUDIT
            ),
            report_type=report_types[3],
        )

        centers_stations_under_invg = stations_and_centers_queryset(
            tally_id=tally_id, qs=qs, report_type=report_types[4]
        )
        centers_stations_ex_after_invg = stations_and_centers_queryset(
            tally_id=tally_id, qs=qs, report_type=report_types[5]
        )

        station_id_query = Subquery(
            Station.objects.filter(
                tally__id=tally_id,
                center__code=OuterRef("center__code"),
                station_number=OuterRef("station_number"),
            ).values("id")[:1],
            output_field=IntegerField(),
        )

        if report_type in [
            "centers-and-stations-in-audit-report",
            "centers-and-stations-under-investigation",
            "centers-and-stations-excluded-after-investigation",
        ]:
            if report_type == "centers-and-stations-in-audit-report":
                self.request.session["station_ids"] = list(
                    centers_stations_in_audit.filter(
                        center__office__region__id=region_id,
                        center__constituency__id=constituency_id,
                    )
                    .annotate(station_id=station_id_query)
                    .values_list("station_id", flat=True)
                )

            if report_type == "centers-and-stations-under-investigation":
                self.request.session["station_ids"] = list(
                    centers_stations_under_invg.filter(
                        center__office__region__id=region_id,
                        center__constituency__id=constituency_id,
                    )
                    .annotate(station_id=station_id_query)
                    .values_list("station_id", flat=True)
                )

            if (
                report_type
                == "centers-and-stations-excluded-after-investigation"
            ):
                self.request.session["station_ids"] = list(
                    centers_stations_ex_after_invg.filter(
                        center__office__region__id=region_id,
                        center__constituency__id=constituency_id,
                    )
                    .annotate(station_id=station_id_query)
                    .values_list("station_id", flat=True)
                )

            return redirect(
                "center-and-stations-list",
                tally_id=tally_id,
                region_id=region_id,
                constituency_id=constituency_id,
            )

        if report_type == "votes-per-candidate-report":
            self.request.session["result_ids"] = list(
                progressive_report.filter(
                    result_form__center__office__region__id=region_id
                ).values_list("id", flat=True)
            )
            self.request.session["ballot_report"] = False

            return redirect(
                "candidate-list-by-votes",
                tally_id=tally_id,
                region_id=region_id,
                constituency_id=constituency_id,
            )

        return self.render_to_response(
            self.get_context_data(
                tally_id=tally_id,
                region_id=region_id,
                administrative_area_name=_("Constituencies"),
                administrative_area_child_report_name=_("Sub Constituencies"),
                turn_out_report_download_url="constituencies-turnout-csv",
                summary_report_download_url="constituencies-summary-csv",
                progressive_report_download_url="constituencies-progressive-csv",
                discrepancy_report_download_url="constituencies-discrepancy-csv",
                centers_stations_under_invg=centers_stations_under_invg,
                centers_stations_ex_after_invg=centers_stations_ex_after_invg,
                region_name=region_name,
                child_summary_report_url="sub-constituency-summary-report",
                child_progressive_report_url="sub-constituency-progressive-report",
                admin_area_votes_per_candidate_report_url="constituency-votes-per-candidate",
                constituency_discrepancy_report_url="constituency-discrepancy-report",
                child_discrepancy_report_url="sub-constituency-discrepancy-report",
                child_admin_area_under_investigation_report_url="sub-constituencies-under-investigation-report",
                child_admin_area_excluded_after_investigation_report_url="sub-constituencies-excluded-after-investigation-report",
            )
        )


class SubConstituencyReportsView(
    LoginRequiredMixin, GroupRequiredMixin, TemplateView
):
    group_required = groups.TALLY_MANAGER

    def get(self, request, *args, **kwargs):
        tally_id = kwargs["tally_id"]
        region_id = kwargs.get("region_id")
        constituency_id = kwargs.get("constituency_id")
        sub_constituency_id = kwargs.get("sub_constituency_id")
        report_type = kwargs.get("report_type")

        region_name = (
            Region.objects.get(id=region_id, tally__id=tally_id).name
            if region_id
            else None
        )
        constituency_name = (
            Constituency.objects.get(
                id=constituency_id, tally__id=tally_id
            ).name
            if constituency_id
            else None
        )

        column_name = "result_form__center__sub_constituency__code"
        progressive_report = generate_progressive_report(
            tally_id=tally_id,
            report_column_name=column_name,
            region_id=region_id,
            constituency_id=constituency_id,
        )

        qs = Station.objects.filter(tally__id=tally_id)

        centers_stations_in_audit = stations_and_centers_queryset(
            tally_id=tally_id,
            qs=ResultForm.objects.filter(
                tally__id=tally_id, form_state=FormState.AUDIT
            ),
            report_type=report_types[3],
        )

        centers_stations_under_invg = stations_and_centers_queryset(
            tally_id=tally_id, qs=qs, report_type=report_types[4]
        )

        centers_stations_ex_after_invg = stations_and_centers_queryset(
            tally_id=tally_id, qs=qs, report_type=report_types[5]
        )

        station_id_query = Subquery(
            Station.objects.filter(
                tally__id=tally_id,
                center__code=OuterRef("center__code"),
                station_number=OuterRef("station_number"),
            ).values("id")[:1],
            output_field=IntegerField(),
        )

        if report_type in [
            "centers-and-stations-in-audit-report",
            "centers-and-stations-under-investigation",
            "centers-and-stations-excluded-after-investigation",
        ]:
            if report_type == "centers-and-stations-in-audit-report":
                self.request.session["station_ids"] = list(
                    centers_stations_in_audit.filter(
                        center__office__region__id=region_id,
                        center__constituency__id=constituency_id,
                        center__sub_constituency__id=sub_constituency_id,
                    )
                    .annotate(station_id=station_id_query)
                    .values_list("station_id", flat=True)
                )

            if report_type == "centers-and-stations-under-investigation":
                self.request.session["station_ids"] = list(
                    centers_stations_under_invg.filter(
                        center__office__region__id=region_id,
                        center__constituency__id=constituency_id,
                        center__sub_constituency__id=sub_constituency_id,
                    )
                    .annotate(station_id=station_id_query)
                    .values_list("station_id", flat=True)
                )

            if (
                report_type
                == "centers-and-stations-excluded-after-investigation"
            ):
                self.request.session["station_ids"] = list(
                    centers_stations_ex_after_invg.filter(
                        center__office__region__id=region_id,
                        center__constituency__id=constituency_id,
                        center__sub_constituency__id=sub_constituency_id,
                    )
                    .annotate(station_id=station_id_query)
                    .values_list("station_id", flat=True)
                )

            return redirect(
                "center-and-stations-list",
                tally_id=tally_id,
                region_id=region_id,
                constituency_id=constituency_id,
                sub_constituency_id=sub_constituency_id,
            )

        if report_type in [
            "votes-per-candidate-report",
            "candidate-list-sorted-by-ballots-number",
        ]:
            self.request.session["result_ids"] = list(
                progressive_report.filter(
                    result_form__center__office__region__id=region_id
                ).values_list("id", flat=True)
            )
            self.request.session["ballot_report"] = (
                report_type == "candidate-list-sorted-by-ballots-number"
            )

            return redirect(
                "candidate-list-by-votes",
                tally_id=tally_id,
                region_id=region_id,
                constituency_id=constituency_id,
                sub_constituency_id=sub_constituency_id,
            )

        return self.render_to_response(
            self.get_context_data(
                tally_id=tally_id,
                region_id=region_id,
                administrative_area_child_report_name=None,
                constituency_id=constituency_id,
                turn_out_report_download_url="sub-constituencies-turnout-csv",
                summary_report_download_url="sub-constituencies-summary-csv",
                progressive_report_download_url="sub-constituencies-progressive-csv",
                admin_area_votes_per_candidate_report_url="sub-constituency-votes-per-candidate",
                discrepancy_report_download_url="sub-constituencies-discrepancy-csv",
                centers_stations_under_invg=centers_stations_under_invg,
                centers_stations_ex_after_invg=centers_stations_ex_after_invg,
                administrative_area_name=_("Sub Constituencies"),
                region_name=region_name,
                constituency_name=constituency_name,
                sub_constituency_discrepancy_report_url="sub-constituency-discrepancy-report",
            )
        )


class ResultFormResultsListDataView(
    LoginRequiredMixin,
    GroupRequiredMixin,
    TallyAccessMixin,
    BaseDatatableView,
):
    group_required = groups.TALLY_MANAGER
    model = Result
    columns = (
        "candidate_name",
        "total_votes",
        "valid_votes",
        "candidate_status",
        "election_level",
        "sub_race_type",
        "order",
        "ballot_number",
    )

    def get(self, request, *args, **kwargs):
        queryset = self.get_initial_queryset()

        total_records = len(queryset)
        page = self.request.POST.get("start", 0)
        page_size = self.request.POST.get("length", 10)
        data = self.request.POST.get("data")
        data = ast.literal_eval(data) if data else None
        if data and parse_int(data.get("export_number")):
            page_size = parse_int(data.get("export_number"))
            total_records = parse_int(data.get("export_number"))

        if page_size == "-1":
            page_records = queryset
        else:
            page_records = queryset[int(page) : int(page) + int(page_size)]

        response_data = JsonResponse(
            {
                "draw": int(self.request.POST.get("draw", 0)),
                "recordsTotal": total_records,
                "recordsFiltered": total_records,
                "data": page_records,
            }
        )

        return response_data

    def get_initial_queryset(self):
        """
        Return the initial queryset for the data table.
        You can modify this method to return a custom list.
        Custom logic to get the data as a list
        """
        tally_id = self.kwargs.get("tally_id")
        data = self.request.POST.get("data")
        keyword = self.request.POST.get("search[value]")

        qs = self.model.objects.select_related(
            "candidate",
        )
        if keyword:
            qs = qs.filter(
                Q(candidate_name__icontains=keyword)
                | Q(
                    candidate__ballot__electrol_race__ballot_name__icontains=keyword
                )
                | Q(
                    candidate__ballot__electrol_race__election_level__icontains=keyword
                )
            )
        qs = qs.filter(
            result_form__tally__id=tally_id,
            result_form__form_state=FormState.ARCHIVED,
            entry_version=EntryVersion.FINAL,
            candidate__ballot__active=True,
            active=True,
        )
        qs = get_filtered_candidate_votes(
            tally_id, qs, data=ast.literal_eval(data) if data else None
        )

        races_ids_to_total_votes_mapping = {}
        for electral_race_id in ElectrolRace.objects.filter(
            tally__id=tally_id
        ).values_list("id", flat=True):
            races_ids_to_total_votes_mapping[electral_race_id] = sum(
                [
                    q.get("total_votes")
                    for q in qs
                    if q.get("electrol_race_id") == electral_race_id
                ]
            )
        results = [
            {
                "candidate_name": result.get("candidate_name"),
                "total_votes": result.get("total_votes"),
                "election_level": result.get("election_level"),
                "sub_race_type": result.get("sub_race_type"),
                "order": result.get("order"),
                "ballot_number": result.get("ballot_number"),
                "candidate_status": result.get("candidate_status"),
                "valid_votes": races_ids_to_total_votes_mapping.get(
                    result.get("electrol_race_id")
                ),
            }
            for result in qs
        ]

        sorted_results_report = sorted(
            results, key=lambda x: -x["total_votes"]
        )

        return sorted_results_report


class ResultFormResultsListView(
    LoginRequiredMixin, GroupRequiredMixin, DataTablesMixin, TemplateView
):
    group_required = groups.TALLY_MANAGER
    model = Result
    template_name = "reports/form_results.html"

    def get(self, request, *args, **kwargs):
        tally_id = kwargs.get("tally_id")
        stations, centers, sub_cons = build_stations_centers_and_sub_cons_list(
            tally_id
        )
        electrol_races = ElectrolRace.objects.filter(tally__id=tally_id)
        columns = (
            "candidate_name",
            "total_votes",
            "valid_votes",
            "candidate_status",
            "election_level",
            "sub_race_type",
            "order",
            "ballot_number",
        )
        dt_columns = [{"data": column} for column in columns]

        ballot_status = [
            {
                "name": "Available For Release",
                "value": "available_for_release",
            },
            {
                "name": "Not Available For Release",
                "value": "not_available_for_release",
            },
        ]
        station_status = [
            {"name": "Active", "value": "active"},
            {"name": "In Active", "value": "inactive"},
        ]
        candidate_status = [
            {"name": "Active", "value": "active"},
            {"name": "In Active", "value": "inactive"},
        ]

        return self.render_to_response(
            self.get_context_data(
                remote_url=reverse("form-results-data", kwargs=kwargs),
                tally_id=tally_id,
                stations=stations,
                centers=centers,
                sub_cons=sub_cons,
                election_level_names=set(
                    electrol_races.values_list("election_level", flat=True)
                ),
                sub_race_type_names=set(
                    electrol_races.values_list("ballot_name", flat=True)
                ),
                ballot_status=ballot_status,
                station_status=station_status,
                candidate_status=candidate_status,
                dt_columns=dt_columns,
                export_url='/ajax/get-export/',
            )
        )


class DuplicateResultsListDataView(
    LoginRequiredMixin,
    GroupRequiredMixin,
    TallyAccessMixin,
    BaseDatatableView,
):
    group_required = groups.TALLY_MANAGER
    model = ResultForm
    columns = (
        "ballot_number",
        "center_code",
        "office",
        "barcode",
        "state",
        "station_number",
        "station_id",
        "votes",
    )

    def filter_queryset(self, qs):
        tally_id = self.kwargs.get("tally_id")
        data = self.request.POST.get("data")
        keyword = self.request.POST.get("search[value]")

        qs = qs.filter(tally__id=tally_id, form_state=FormState.UNSUBMITTED)

        duplicate_result_forms = get_result_form_with_duplicate_results(
            tally_id=tally_id, qs=qs
        )

        qs = (
            duplicate_results_queryset(
                tally_id=tally_id,
                qs=duplicate_result_forms,
                data=ast.literal_eval(data),
            )
            if data
            else duplicate_results_queryset(
                tally_id=tally_id, qs=duplicate_result_forms
            )
        )

        if keyword:
            qs = qs.filter(
                Q(votes__contains=keyword)
                | Q(barcode__icontains=keyword)
                | Q(ballot_number__contains=keyword)
                | Q(station_id__contains=keyword)
                | Q(center_code__contains=keyword)
                | Q(state__icontains=keyword)
            )
        return qs

    def render_column(self, row, column):
        if column == "ballot_number":
            return str(f'<td class="center">{row["ballot_number"]}</td>')
        elif column == "center_code":
            return str(f'<td class="center">{row["center_code"]}</td>')
        elif column == "barcode":
            return str(f'<td class="center">{row["barcode"]}</td>')
        elif column == "state":
            return str(f'<td class="center">{row["state"].name}</td>')
        elif column == "station_number":
            return str(f'<td class="center">{row["station_number"]}</td>')
        elif column == "station_id":
            return str(f'<td class="center">{row["station_id"]}</td>')
        elif column == "votes":
            return str(f'<td class="center">{row["votes"]}</td>')
        else:
            return super(DuplicateResultsListDataView, self).render_column(
                row, column
            )


class DuplicateResultsListView(
    LoginRequiredMixin, GroupRequiredMixin, DataTablesMixin, TemplateView
):
    group_required = groups.TALLY_MANAGER
    model = ResultForm
    template_name = "reports/duplicate_results.html"

    def get(self, request, *args, **kwargs):
        tally_id = kwargs.get("tally_id")
        stations, centers, _unused = build_stations_centers_and_sub_cons_list(
            tally_id
        )

        return self.render_to_response(
            self.get_context_data(
                remote_url=reverse("duplicate-results-data", kwargs=kwargs),
                tally_id=tally_id,
                stations=stations,
                centers=centers,
            )
        )


class AllCandidatesVotesDataView(
    LoginRequiredMixin,
    GroupRequiredMixin,
    TallyAccessMixin,
    BaseDatatableView,
):
    group_required = groups.TALLY_MANAGER
    model = AllCandidatesVotes
    columns = (
        "full_name",
        "total_votes",
        "candidate_votes_included_quarantine",
        "stations",
        "stations_completed",
        "stations_complete_percent",
        "ballot_number",
    )

    def filter_queryset(self, qs):
        tally_id = self.kwargs.get("tally_id")
        data = self.request.POST.get("data")
        keyword = self.request.POST.get("search[value]")
        qs = qs.values(
            "ballot_number",
            "stations",
            "stations_completed",
            "stations_complete_percent",
            "full_name",
            "total_votes",
            "candidate_votes_included_quarantine",
            "center_ids",
            "station_numbers",
        ).filter(tally_id=tally_id)

        if data:
            qs = filter_candidates_votes_queryset(
                qs=qs, data=ast.literal_eval(data)
            )

        if keyword:
            qs = qs.filter(
                Q(full_name__icontains=keyword)
                | Q(total_votes__contains=keyword)
                | Q(candidate_votes_included_quarantine__contains=keyword)
                | Q(stations_completed__contains=keyword)
                | Q(ballot_number__contains=keyword)
            )

        return qs

    def render_column(self, row, column):
        if column == "full_name":
            return str(f'<td class="center">{row["full_name"]}</td>')
        elif column == "total_votes":
            return str(f'<td class="center">{row["total_votes"]}</td>')
        elif column == "candidate_votes_included_quarantine":
            return str(
                '<td class="center">'
                f"{row['candidate_votes_included_quarantine']}</td>"
            )
        elif column == "stations":
            return str(f'<td class="center">{row["stations"]}</td>')
        elif column == "stations_completed":
            return str(f'<td class="center">{row["stations_completed"]}</td>')
        elif column == "stations_complete_percent":
            return str(
                f'<td class="center">{row["stations_complete_percent"]}</td>'
            )
        elif column == "ballot_number":
            return str(f'<td class="center">{row["ballot_number"]}</td>')
        else:
            return super(AllCandidatesVotesDataView, self).render_column(
                row, column
            )


class AllCandidatesVotesListView(
    LoginRequiredMixin, GroupRequiredMixin, DataTablesMixin, TemplateView
):
    group_required = groups.TALLY_MANAGER
    model = Station
    template_name = "reports/candidates_votes_report.html"

    def get(self, request, *args, **kwargs):
        tally_id = kwargs.get("tally_id")
        stations, centers, _unused = build_stations_centers_and_sub_cons_list(
            tally_id
        )

        return self.render_to_response(
            self.get_context_data(
                remote_url=reverse("all-candidates-votes-data", kwargs=kwargs),
                tally_id=tally_id,
                stations=stations,
                centers=centers,
                title=_("All Candidates Votes"),
                export_file_name=_("all_candidates_votes"),
            )
        )


class ActiveCandidatesVotesDataView(
    LoginRequiredMixin,
    GroupRequiredMixin,
    TallyAccessMixin,
    BaseDatatableView,
):
    group_required = groups.TALLY_MANAGER
    model = AllCandidatesVotes
    columns = (
        "full_name",
        "total_votes",
        "candidate_votes_included_quarantine",
        "stations",
        "stations_completed",
        "stations_complete_percent",
        "ballot_number",
    )

    def filter_queryset(self, qs):
        tally_id = self.kwargs.get("tally_id")
        data = self.request.POST.get("data")
        keyword = self.request.POST.get("search[value]")

        qs = qs.values(
            "ballot_number",
            "stations",
            "stations_completed",
            "stations_complete_percent",
            "full_name",
            "total_votes",
            "candidate_votes_included_quarantine",
            "center_ids",
            "station_numbers",
        ).filter(tally_id=tally_id, candidate_active=True)

        if data:
            qs = filter_candidates_votes_queryset(
                qs=qs, data=ast.literal_eval(data)
            )

        if keyword:
            qs = qs.filter(
                Q(full_name__icontains=keyword)
                | Q(total_votes__contains=keyword)
                | Q(candidate_votes_included_quarantine__contains=keyword)
                | Q(stations_completed__contains=keyword)
                | Q(ballot_number__contains=keyword)
            )

        return qs

    def render_column(self, row, column):
        if column == "full_name":
            return str(f'<td class="center">{row["full_name"]}</td>')
        elif column == "total_votes":
            return str(f'<td class="center">{row["total_votes"]}</td>')
        elif column == "candidate_votes_included_quarantine":
            return str(
                '<td class="center">'
                f"{row['candidate_votes_included_quarantine']}</td>"
            )
        elif column == "stations":
            return str(f'<td class="center">{row["stations"]}</td>')
        elif column == "stations_completed":
            return str(f'<td class="center">{row["stations_completed"]}</td>')
        elif column == "stations_complete_percent":
            return str(
                f'<td class="center">{row["stations_complete_percent"]}</td>'
            )
        elif column == "ballot_number":
            return str(f'<td class="center">{row["ballot_number"]}</td>')
        else:
            return super(ActiveCandidatesVotesDataView, self).render_column(
                row, column
            )


class ActiveCandidatesVotesListView(
    LoginRequiredMixin, GroupRequiredMixin, DataTablesMixin, TemplateView
):
    group_required = groups.TALLY_MANAGER
    model = Station
    template_name = "reports/candidates_votes_report.html"

    def get(self, request, *args, **kwargs):
        tally_id = kwargs.get("tally_id")
        stations, centers, _unused = build_stations_centers_and_sub_cons_list(
            tally_id
        )

        return self.render_to_response(
            self.get_context_data(
                remote_url=reverse(
                    "active-candidates-votes-data", kwargs=kwargs
                ),
                tally_id=tally_id,
                stations=stations,
                centers=centers,
                title=_("Active Candidates Votes"),
                export_file_name=_("active_candidates_votes"),
            )
        )


class ClearanceAuditSummaryReportView(
    LoginRequiredMixin,
    GroupRequiredMixin,
    DataTablesMixin,
    TemplateView,
):
    group_required = [groups.SUPER_ADMINISTRATOR, groups.TALLY_MANAGER]
    template_name = "reports/clearance_audit_summary.html"

    def get(self, request, *args, **kwargs):
        tally_id = self.kwargs.get("tally_id")
        columns = (
            "barcode",
            "center_code",
            "station_number",
            "race",
            "sub_race",
            "municipality_name",
            "municipality_code",
            "action_prior",
            "recommendation",
            "decision",
            "issue_reason",
            "supervisor",
            "last_modified",
        )
        dt_columns = [{"data": column} for column in columns]
        _, _, sub_cons = build_stations_centers_and_sub_cons_list(tally_id)
        electrol_races = ElectrolRace.objects.filter(tally__id=tally_id)
        context_data = {
            "tally_id": tally_id,
            "remote_url": reverse(
                "clearance-audit-summary-data",
                kwargs={"tally_id": tally_id},
            ),
            "sub_cons": sub_cons,
            "races": set(
                electrol_races.values_list("election_level", flat=True)
            ),
            "sub_races": set(
                electrol_races.values_list("ballot_name", flat=True)
            ),
            "dt_columns": dt_columns,
        }
        return self.render_to_response(self.get_context_data(**context_data))


class ClearanceAuditSummaryReportDataView(
    LoginRequiredMixin,
    GroupRequiredMixin,
    TallyAccessMixin,
    BaseDatatableView,
):
    group_required = [groups.SUPER_ADMINISTRATOR, groups.TALLY_MANAGER]
    columns = [
        "barcode",
        "center_code",
        "station_number",
        "race",
        "sub_race",
        "municipality_name",
        "municipality_code",
        "action_prior",
        "recommendation",
        "decision",
        "issue_reason",
        "supervisor",
        "last_modified",
    ]
    order_columns = columns

    def get_initial_queryset(self, data=None):
        tally_id = self.kwargs.get("tally_id")
        qs = ResultForm.objects.filter(tally_id=tally_id)
        tab = "clearance"
        municipalities = None
        races = None
        sub_races = None
        if data:
            tab = data.get("tab", "clearance")
            municipalities = data.get("sub_con_codes")
            races = data.get("election_level_names")
            sub_races = data.get("sub_race_type_names")
        if tab == "clearance":
            qs = qs.filter(clearances__active=True)
        elif tab == "audit":
            qs = qs.filter(form_state=FormState.AUDIT)
        if municipalities:
            qs = qs.filter(center__sub_constituency__code__in=municipalities)
        if races:
            qs = qs.filter(ballot__electrol_race__election_level__in=races)
        if sub_races:
            qs = qs.filter(ballot__electrol_race__ballot_name__in=sub_races)
        qs = qs.select_related(
            "center",
            "ballot",
            "center__office",
            "ballot__electrol_race",
            "center__sub_constituency",
        )
        return [
            {
                "barcode": form.barcode,
                "center_code": form.center.code if form.center else "",
                "station_number": form.station_number,
                "race": (
                    form.ballot.electrol_race.election_level
                    if form.ballot and form.ballot.electrol_race
                    else ""
                ),
                "sub_race": (
                    form.ballot.electrol_race.ballot_name
                    if form.ballot and form.ballot.electrol_race
                    else ""
                ),
                "municipality_name": (
                    form.center.sub_constituency.name
                    if form.center and form.center.sub_constituency
                    else ""
                ),
                "municipality_code": (
                    form.center.sub_constituency.code
                    if form.center and form.center.sub_constituency
                    else ""
                ),
                "action_prior": (
                    getattr(form, "clearance_action_prior", "")
                    if tab == "clearance"
                    else getattr(form, "audit_action_prior", "")
                ),
                "recommendation": (
                    getattr(form, "clearance_recommendation", "")
                    if tab == "clearance"
                    else getattr(form, "audit_recommendation", "")
                ),
                "decision": getattr(form, "decision", ""),
                "issue_reason": getattr(form, "issue_reason", ""),
                "supervisor": getattr(form, "supervisor", ""),
                "last_modified": getattr(form, "modified_date_formatted", ""),
            }
            for form in qs
        ]

    def get(self, request, *args, **kwargs):
        request_data = request.GET.get("data")
        data = None
        if request_data:
            try:
                data = ast.literal_eval(request_data)
            except Exception:
                data = None
        queryset = self.get_initial_queryset(data)
        total_records = len(queryset)
        page = request.GET.get("start", 0)
        page_size = request.GET.get("length", 10)
        search = request.GET.get("search[value]", None)

        if search:
            queryset = [
                row for row in queryset if search.lower() in str(row).lower()
            ]
            total_records = len(queryset)

        if page_size == "-1":
            page_records = queryset
        else:
            page_records = queryset[int(page) : int(page) + int(page_size)]

        response_data = JsonResponse(
            {
                "draw": int(request.GET.get("draw", 0)),
                "recordsTotal": total_records,
                "recordsFiltered": total_records,
                "data": page_records,
            }
        )
        return response_data
